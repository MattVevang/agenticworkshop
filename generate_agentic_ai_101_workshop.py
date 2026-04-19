#!/usr/bin/env python3
"""Generate an AI 101 + Agentic Tools workshop deck (.pptx).

Usage:
    pip install python-pptx
    python generate_agentic_ai_101_workshop.py

Output:
    Agentic_AI_101_Workshop.pptx (next to this script)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme: "Neon Mission Control"
# ---------------------------------------------------------------------------
BG = RGBColor(0x0B, 0x10, 0x1A)            # near-black blue
SURFACE = RGBColor(0x13, 0x1B, 0x2E)       # panel background
PRIMARY = RGBColor(0x22, 0xD3, 0xEE)       # cyan accent
SECONDARY = RGBColor(0xA7, 0x8B, 0xFA)     # purple accent
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)       # green
WARN = RGBColor(0xF5, 0x9E, 0x0B)          # amber
DANGER = RGBColor(0xEF, 0x44, 0x44)        # red
TEXT = RGBColor(0xE5, 0xE7, 0xEB)          # light gray
MUTED = RGBColor(0x9C, 0xA3, 0xAF)         # muted gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def override_last_slide_notes(prs, text):
    """Replace notes on the most recently added slide."""
    slide = prs.slides[len(prs.slides) - 1]
    slide.notes_slide.notes_text_frame.text = text


def add_background(slide, width, height):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    # Accent blobs for a more visual theme.
    orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(-0.8), Inches(3.2), Inches(3.2))
    orb1.fill.solid()
    orb1.fill.fore_color.rgb = PRIMARY
    orb1.fill.transparency = 0.82
    orb1.line.fill.background()

    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.0), Inches(5.2), Inches(3.8), Inches(3.8))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = SECONDARY
    orb2.fill.transparency = 0.86
    orb2.line.fill.background()


def add_header(slide, width, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SURFACE
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), Inches(0.95))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(11.8), Inches(0.66))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(0)
        p2.space_after = Pt(0)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(12)
        run2.font.color.rgb = MUTED


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.2), Inches(0.3))
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, fill=SURFACE, line=PRIMARY):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = line
    panel.line.width = Pt(1.3)
    return panel


def set_first_bullet(tf, text, size=Pt(18), bold=False, color=TEXT):
    p = tf.paragraphs[0]
    p.level = 0
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_bullet(tf, text, level=0, size=Pt(16), bold=False, color=TEXT):
    p = tf.add_paragraph()
    # Keep PowerPoint-level indentation flat to avoid layout drift on blank slides.
    # Render nested structure with textual indent instead.
    p.level = 0
    p.space_after = Pt(4)
    run = p.add_run()
    prefix = ("    " * level) if level > 0 else ""
    run.text = f"{prefix}{text}"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_title_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.8), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI 101 + Agentic Tools Workshop"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "From model basics to practical workflows for students"
    run2.font.size = Pt(24)
    run2.font.color.rgb = PRIMARY

    add_footer(slide, "Agentic Workshop Lab")
    add_notes(
        slide,
        "SPEAKER NOTES — Title Slide\n"
        "• Welcome students. Introduce yourself and any co-presenters.\n"
        "• Set expectations: this is hands-on, interactive, and fun — not a lecture.\n"
        "• Mention: 'By the end of today you'll have chatted with real AI models running on hardware in this room.'\n"
        "• Transition: 'Let's start with what we're going to cover today.'\n"
        "• Timing: ~2 min\n\n"
        "Image seed: 'cinematic neon classroom mission control, teenagers learning AI, high contrast, wide shot'",
    )


def add_bullets_slide(prs, blank, title, bullets, image_hint):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, title)

    panel = add_panel(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.55), fill=SURFACE, line=PRIMARY)
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.18)

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            head, subs = item
            if i == 0:
                set_first_bullet(tf, head, size=Pt(20), bold=True, color=WHITE)
            else:
                add_bullet(tf, head, size=Pt(20), bold=True, color=WHITE)
            for sub in subs:
                add_bullet(tf, sub, level=1, size=Pt(14), color=TEXT)
        else:
            if i == 0:
                set_first_bullet(tf, item, size=Pt(18), color=TEXT)
            else:
                add_bullet(tf, item, size=Pt(18), color=TEXT)

    add_footer(slide, "AI 101 Workshop")
    add_notes(slide, f"Image concept: {image_hint}")


def add_two_column_slide(prs, blank, title, left_title, left_points, right_title, right_points, image_hint):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, title)

    left = add_panel(slide, Inches(0.6), Inches(1.25), Inches(5.85), Inches(5.35), fill=SURFACE, line=PRIMARY)
    right = add_panel(slide, Inches(6.85), Inches(1.25), Inches(5.85), Inches(5.35), fill=SURFACE, line=SECONDARY)

    tf_l = left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.2)
    tf_l.margin_top = Inches(0.1)
    set_first_bullet(tf_l, left_title, size=Pt(21), bold=True, color=PRIMARY)
    for point in left_points:
        add_bullet(tf_l, point, size=Pt(14), color=TEXT)

    tf_r = right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.2)
    tf_r.margin_top = Inches(0.1)
    set_first_bullet(tf_r, right_title, size=Pt(21), bold=True, color=SECONDARY)
    for point in right_points:
        add_bullet(tf_r, point, size=Pt(14), color=TEXT)

    add_footer(slide, "AI 101 Workshop")
    add_notes(slide, f"Image concept: {image_hint}\n(Add speaker talking points for this slide in the generator's build_deck() function.)")


def add_model_size_metaphor_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "Model Scale, Parameters, and Tradeoffs")

    explainer = add_panel(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.35), fill=SURFACE, line=PRIMARY)
    tf = explainer.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.07)
    set_first_bullet(tf, "Parameter count is not IQ, but it affects capability range, latency, and cost.", size=Pt(17), bold=True, color=WHITE)
    add_bullet(tf, "Bigger models usually reason deeper and transfer better. Smaller models are faster and cheaper.", size=Pt(13), color=TEXT)

    # Visual metaphor cards
    big_card = add_panel(slide, Inches(0.8), Inches(2.85), Inches(5.8), Inches(3.35), fill=SURFACE, line=SUCCESS)
    tfb = big_card.text_frame
    tfb.word_wrap = True
    tfb.margin_left = Inches(0.2)
    tfb.margin_top = Inches(0.1)
    set_first_bullet(tfb, "Think-Book Model", size=Pt(24), bold=True, color=SUCCESS)
    add_bullet(tfb, "Deep context and nuance", size=Pt(14), color=TEXT)
    add_bullet(tfb, "Great for architecture, debugging, and tradeoffs", size=Pt(14), color=TEXT)
    add_bullet(tfb, "Higher cost + longer response time", size=Pt(14), color=TEXT)

    small_card = add_panel(slide, Inches(6.75), Inches(2.85), Inches(5.8), Inches(3.35), fill=SURFACE, line=WARN)
    tfs = small_card.text_frame
    tfs.word_wrap = True
    tfs.margin_left = Inches(0.2)
    tfs.margin_top = Inches(0.1)
    set_first_bullet(tfs, "Cliff-Notes Model", size=Pt(24), bold=True, color=WARN)
    add_bullet(tfs, "Quick, cheap, and responsive", size=Pt(14), color=TEXT)
    add_bullet(tfs, "Great for summaries, rewrites, lightweight tasks", size=Pt(14), color=TEXT)
    add_bullet(tfs, "Can miss subtle constraints in hard code tasks", size=Pt(14), color=TEXT)

    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — Model Scale, Parameters, and Tradeoffs\n"
        "• Key point: bigger isn't always better — it depends on the task.\n"
        "• Ask the audience: 'Would you use a sledgehammer to hang a picture frame?'\n"
        "• Explain: a 1B model is like cliff notes — fast but shallow. A 70B model is like a reference textbook.\n"
        "• Connect to their experience: 'The models you'll use today range from 1B to 14B parameters.'\n"
        "• Transition: 'Let's see what hardware actually runs these models.'\n"
        "• Timing: ~3 min\n\n"
        "Image concept: side-by-side desk with a thick reference textbook vs slim cheat-sheet notebook.",
    )


def add_tools_landscape_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "Tool Landscape: Copilot vs Claude vs OpenCode-style Agents")

    cols = [
        ("GitHub Copilot", PRIMARY, [
            "Best inside developer workflow and repos",
            "Strong for inline coding, PR workflows, and agent handoff",
            "Great when GitHub context matters most",
        ]),
        ("Claude-style Assistant", SECONDARY, [
            "Strong long-form reasoning and explanation",
            "Useful for planning, docs, and architecture tradeoffs",
            "Great for teaching and concept walkthroughs",
        ]),
        ("OpenCode / DIY Agents", SUCCESS, [
            "Maximum control over models, tools, and infra",
            "Good for experimentation and custom orchestration",
            "Higher setup and governance overhead",
        ]),
    ]

    x = Inches(0.55)
    for title, color, points in cols:
        card = add_panel(slide, x, Inches(1.35), Inches(4.1), Inches(5.2), fill=SURFACE, line=color)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.08)
        set_first_bullet(tf, title, size=Pt(20), bold=True, color=color)
        for p in points:
            add_bullet(tf, p, size=Pt(13), color=TEXT)
        x += Inches(4.28)

    note = add_panel(slide, Inches(0.55), Inches(6.65), Inches(12.15), Inches(0.6), fill=BG, line=PRIMARY)
    nt = note.text_frame
    nt.margin_left = Inches(0.15)
    nt.margin_top = Inches(0.03)
    set_first_bullet(nt, "Pick tools by workflow fit, not by hype. Blended stacks often win.", size=Pt(12), bold=True, color=WHITE)
    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — Tool Landscape\n"
        "• This is NOT a 'which is best' slide — emphasize workflow fit.\n"
        "• Ask: 'How many of you have heard of GitHub Copilot? ChatGPT? Any others?'\n"
        "• Key point: different tools solve different problems. A hammer vs screwdriver vs wrench.\n"
        "• Copilot = best inside code. Claude = best for long thinking. Open tools = full control.\n"
        "• Mention: 'Blended stacks often win — you might use multiple tools in a day.'\n"
        "• Transition: 'Now let's look at how agents connect to the real world...'\n"
        "• Timing: ~3 min\n\n"
        "Image concept: toolbelt metaphor with 3 labeled tools (integrated, reasoning-focused, customizable).",
    )


def add_mcp_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "What is an MCP Server?")

    center = add_panel(slide, Inches(4.95), Inches(2.5), Inches(3.4), Inches(1.5), fill=SURFACE, line=PRIMARY)
    tcenter = center.text_frame
    tcenter.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tcenter.paragraphs[0].add_run()
    run.text = "AI Agent"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE

    spokes = [
        (Inches(0.9), Inches(1.5), "Filesystem\nTools"),
        (Inches(0.9), Inches(4.4), "Browser\nAutomation"),
        (Inches(9.25), Inches(1.5), "Tickets +\nDocs APIs"),
        (Inches(9.25), Inches(4.4), "Monitoring /\nTelemetry"),
    ]
    for x, y, label in spokes:
        box = add_panel(slide, x, y, Inches(3.1), Inches(1.35), fill=SURFACE, line=SECONDARY)
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = TEXT

        connector = slide.shapes.add_shape(
            MSO_SHAPE.LINE_INVERSE,
            x + Inches(1.55),
            y + Inches(0.68),
            Inches(3.35),
            Inches(1.3),
        )
        connector.line.color.rgb = PRIMARY
        connector.line.width = Pt(1.2)

    bottom = add_panel(slide, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.9), fill=SURFACE, line=SUCCESS)
    tfb = bottom.text_frame
    tfb.margin_left = Inches(0.2)
    tfb.margin_top = Inches(0.08)
    set_first_bullet(tfb, "MCP is a standard way to expose tools safely so agents can do real work, not just chat.", size=Pt(14), bold=True, color=SUCCESS)
    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — MCP (Model Context Protocol)\n"
        "• Start with the problem: 'AI can chat, but what if you want it to actually DO things?'\n"
        "• Explain: MCP is like USB-C for AI — one standard plug for many tools.\n"
        "• Walk through the diagram: the agent in the center can talk to files, browsers, APIs, monitoring.\n"
        "• Real example: 'An AI agent could search your codebase, run tests, and create a pull request.'\n"
        "• Ask: 'What tools would YOU want to give an AI agent access to?'\n"
        "• Transition: 'Let's talk about the currency of AI — tokens.'\n"
        "• Timing: ~3 min\n\n"
        "Image concept: hub-and-spoke system diagram with AI agent in center and tool systems around it.",
    )


def add_tokens_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "Tokens, Context Windows, and Limits")

    left = add_panel(slide, Inches(0.6), Inches(1.25), Inches(6.0), Inches(5.45), fill=SURFACE, line=PRIMARY)
    tf_l = left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.2)
    tf_l.margin_top = Inches(0.08)
    set_first_bullet(tf_l, "Token basics", size=Pt(22), bold=True, color=PRIMARY)
    add_bullet(tf_l, "A token is a chunk of text the model processes.", size=Pt(14))
    add_bullet(tf_l, "Input tokens + output tokens drive cost and limits.", size=Pt(14))
    add_bullet(tf_l, "Context window is the active working memory budget.", size=Pt(14))
    add_bullet(tf_l, "Long chats can force summarization/compression.", size=Pt(14))

    right = add_panel(slide, Inches(6.75), Inches(1.25), Inches(6.0), Inches(5.45), fill=SURFACE, line=DANGER)
    tf_r = right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.2)
    tf_r.margin_top = Inches(0.08)
    set_first_bullet(tf_r, "When you exceed context", size=Pt(22), bold=True, color=DANGER)
    add_bullet(tf_r, "Older details get dropped or compressed.", size=Pt(14))
    add_bullet(tf_r, "Precision can degrade on long, complex tasks.", size=Pt(14))
    add_bullet(tf_r, "Symptoms: repeats, contradictions, missing constraints.", size=Pt(14))
    add_bullet(tf_r, "Fixes: restate constraints, chunk tasks, refresh summary.", size=Pt(14))

    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — Tokens, Context Windows, and Limits\n"
        "• Explain: tokens are the 'currency' of AI. Every word costs tokens. There's a budget.\n"
        "• Analogy: 'Imagine a desk that can only hold 10 pages. Add page 11, page 1 falls off.'\n"
        "• Left panel: how tokens work. Right panel: what happens when you run out.\n"
        "• Practical tip: 'This is why long chats get weird — the model literally forgets the beginning.'\n"
        "• Ask: 'Has anyone had a chatbot seem to forget what you told it earlier?'\n"
        "• Transition: 'Let's compare parameters vs tokens — they sound similar but are very different.'\n"
        "• Timing: ~3 min\n\n"
        "Image concept: backpack with finite capacity; adding more items forces throwing old ones out.",
    )


def add_storyboard_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "Imagery Storyboard Suggestions")

    prompts = [
        ("Model size metaphor", "thick reference book vs slim cliff-notes booklet on a desk, split-frame, high clarity"),
        ("Local vs cloud", "laptop with small on-device chip vs cloud datacenter pipeline, clean infographic style"),
        ("MCP concept", "AI hub connected to tools: browser, docs, terminal, monitoring, network graph style"),
        ("Context overflow", "container filling to limit and compressing old notes into a summary card"),
        ("Right model for task", "task board mapping quick tasks to small model and deep architecture to large model"),
    ]

    table_shape = slide.shapes.add_table(1 + len(prompts), 2, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.9))
    table = table_shape.table
    table.columns[0].width = Inches(3.1)
    table.columns[1].width = Inches(9.0)

    headers = ["Slide Idea", "Prompt Seed (edit for your image generator)"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE

    for r, (label, prompt) in enumerate(prompts, start=1):
        c0 = table.cell(r, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = BG
        c0.text = label
        p0 = c0.text_frame.paragraphs[0]
        p0.runs[0].font.size = Pt(11)
        p0.runs[0].font.bold = True
        p0.runs[0].font.color.rgb = PRIMARY

        c1 = table.cell(r, 1)
        c1.fill.solid()
        c1.fill.fore_color.rgb = BG
        c1.text = prompt
        p1 = c1.text_frame.paragraphs[0]
        p1.runs[0].font.size = Pt(11)
        p1.runs[0].font.color.rgb = TEXT

    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — Imagery Storyboard\n"
        "• This slide is a reference for the presenter — skip in live presentation if time is tight.\n"
        "• Use these prompt seeds to generate custom images for the deck in Midjourney, DALL-E, or Stable Diffusion.\n"
        "• Alternatively, search stock photo sites using the slide idea as a keyword.\n"
        "• Timing: skip or ~1 min",
    )


# ---------------------------------------------------------------------------
# New topic slides
# ---------------------------------------------------------------------------


def add_training_deep_dive_slide(prs, blank):
    """Slide: How Training Really Works (deeper dive)."""
    add_bullets_slide(
        prs, blank, "How Training Really Works",
        [
            ("Data Collection", [
                "Internet text, books, code repos \u2014 trillions of words scraped and cleaned."
            ]),
            ("Pretraining", [
                "Model learns to predict the next token across massive datasets.",
                "Like reading millions of books \u2014 you absorb grammar and style."
            ]),
            ("Fine-Tuning with Human Feedback (RLHF)", [
                "Humans rate outputs to steer the model toward helpful, safe responses."
            ]),
            ("Key Insight", [
                "Models learn patterns, not facts \u2014 they are pattern engines.",
                "You might misremember specific facts, but you \"get\" the language."
            ]),
        ],
        "data pipeline: books and code flowing through neural network into a glowing brain",
    )


def add_prediction_slide(prs, blank):
    """Slide: How Does Prediction Actually Work?"""
    add_bullets_slide(
        prs, blank, "How Does Prediction Actually Work?",
        [
            ("Token-by-Token Generation", [
                "Model picks the most likely next token, feeds it back, repeats.",
                "Like autocomplete on steroids \u2014 instead of one word, it keeps going."
            ]),
            ("Temperature Controls Randomness", [
                "Low temperature = focused, predictable output.",
                "High temperature = more creative, surprising results."
            ]),
            ("Top-K and Top-P Sampling", [
                "Only consider the top candidates for the next token.",
                "Balances quality and variety in responses."
            ]),
            "Live Demo: We have a token prediction demo that shows this in action!",
        ],
        "autocomplete cascade showing probability bars for next-token candidates",
    )


def add_hallucinations_slide(prs, blank):
    """Slide: Hallucinations \u2014 When AI Makes Stuff Up."""
    add_bullets_slide(
        prs, blank, "Hallucinations: When AI Makes Stuff Up",
        [
            ("What Happens", [
                "Model generates confident-sounding but factually wrong output."
            ]),
            ("Why It Happens", [
                "It predicts likely tokens, not verified facts \u2014 no fact-checker inside."
            ]),
            ("Examples You Might See", [
                "Fake book summaries, invented citations, wrong math steps."
            ]),
            ("How to Minimize", [
                "Be specific, ask for sources, fact-check, use smaller-scope questions."
            ]),
            "Trust but verify \u2014 treat AI like a smart but unreliable first draft.",
        ],
        "magnifying glass revealing cracks in a confident-looking AI-generated document",
    )


def add_multimodal_slide(prs, blank):
    """Slide: Why Can Some Models Generate Images?"""
    add_bullets_slide(
        prs, blank, "Why Can Some Models Generate Images?",
        [
            ("Text Models", [
                "Trained on text data only \u2014 read and write text, nothing else."
            ]),
            ("Multimodal Models (GPT-4o, Gemini)", [
                "Trained on text + image pairs \u2014 can understand photos you upload."
            ]),
            ("Image Generation (Stable Diffusion, DALL-E)", [
                "Totally different architecture \u2014 uses a diffusion process.",
                "Starts from noise, gradually refines into an image."
            ]),
            ("Key Insight", [
                "A model can only do what it was trained to do.",
                "Text model \u2260 image model \u2014 different tools for different jobs."
            ]),
        ],
        "split view: text tokens flowing through one model, pixel grids through another",
    )


def add_hardware_slide(prs, blank):
    """Slide: What Hardware Runs a Model?"""
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "What Hardware Runs a Model?")

    explainer = add_panel(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.35),
                          fill=SURFACE, line=PRIMARY)
    tf = explainer.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.07)
    set_first_bullet(tf, "GPUs power AI because they do massive parallel math. VRAM is the bottleneck.",
                     size=Pt(17), bold=True, color=WHITE)
    add_bullet(tf, "CPU = general-purpose brain.  GPU = thousands of simple cores doing math at once.",
               size=Pt(13), color=TEXT)

    rows_data = [
        ("1B params", "~2 GB", "Tiny tasks, mobile devices"),
        ("7B params", "~6 GB", "Good starter, fits most gaming GPUs"),
        ("14B params", "~10 GB", "Solid quality, needs mid-range GPU"),
        ("30B params", "~20 GB", "Near-cloud quality, high-end GPU"),
        ("70B+ params", "~48 GB+", "Top tier \u2014 multiple GPUs needed"),
    ]

    table_shape = slide.shapes.add_table(
        1 + len(rows_data), 3, Inches(0.8), Inches(2.85), Inches(11.7), Inches(3.0))
    table = table_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(5.7)

    for c, text in enumerate(["Model Size", "VRAM Needed", "What That Means"]):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE

    for r, (size, vram, note) in enumerate(rows_data, start=1):
        for c, val in enumerate([size, vram, note]):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = TEXT
            if c == 0:
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = PRIMARY

    callout = add_panel(slide, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.7),
                        fill=SURFACE, line=SUCCESS)
    tfc = callout.text_frame
    tfc.margin_left = Inches(0.2)
    tfc.margin_top = Inches(0.06)
    set_first_bullet(
        tfc,
        "Our workshop server: RTX 5090 with 32 GB VRAM \u2014 runs up to ~30B models comfortably.",
        size=Pt(14), bold=True, color=SUCCESS)

    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — What Hardware Runs a Model?\n"
        "• Key message: AI needs special hardware — GPUs, not just CPUs.\n"
        "• Walk through the table: 'The tiny 1B model fits anywhere. The 70B model needs serious hardware.'\n"
        "• Point to the callout: 'Our server has an RTX 5090 with 32GB VRAM — that's why we can run 14B models.'\n"
        "• Ask: 'How many of you have a gaming GPU at home? You might be able to run AI on it!'\n"
        "• Fun fact: 'The GPU in a PS5 could run a small AI model.'\n"
        "• Transition: 'Now that you know the hardware, let's understand the difference between parameters and tokens.'\n"
        "• Timing: ~3 min\n\n"
        "Image concept: GPU card with glowing VRAM chips, size comparison chart.",
    )


def add_params_vs_tokens_slide(prs, blank):
    """Slide: Parameters vs Tokens \u2014 What's the Difference?"""
    add_two_column_slide(
        prs, blank, "Parameters vs Tokens: What's the Difference?",
        "Parameters (The Engine)",
        [
            "The model's learned knowledge \u2014 its neural weights.",
            "Fixed after training. Like the brain's connections.",
            "Parameter count = model size and capability.",
            "More parameters = deeper reasoning potential.",
            "Analogy: the size of the engine under the hood.",
        ],
        "Tokens (The Fuel)",
        [
            "Chunks of text flowing through at runtime.",
            "Like the words in a conversation \u2014 always changing.",
            "Token count = conversation length and cost.",
            "More tokens = longer context, higher price tag.",
            "Analogy: the fuel going through the engine.",
        ],
        "engine-vs-fuel metaphor: car engine cutaway on left, fuel gauge on right",
    )


def add_privacy_slide(prs, blank):
    """Slide: Privacy and Your Data \u2014 two columns with warning callout."""
    slide = prs.slides.add_slide(blank)
    add_background(slide, prs.slide_width, prs.slide_height)
    add_header(slide, prs.slide_width, "Privacy and Your Data")

    left = add_panel(slide, Inches(0.6), Inches(1.25), Inches(5.85), Inches(4.6),
                     fill=SURFACE, line=PRIMARY)
    right = add_panel(slide, Inches(6.85), Inches(1.25), Inches(5.85), Inches(4.6),
                      fill=SURFACE, line=SECONDARY)

    tf_l = left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.2)
    tf_l.margin_top = Inches(0.1)
    set_first_bullet(tf_l, "Cloud AI Services", size=Pt(21), bold=True, color=PRIMARY)
    for pt in [
        "Your prompts may be logged by the provider.",
        "Some services use your data for training (check!).",
        "Subject to the provider's data policies.",
        "Consider carefully what you share.",
        "Great capability, but read the fine print.",
    ]:
        add_bullet(tf_l, pt, size=Pt(14), color=TEXT)

    tf_r = right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.2)
    tf_r.margin_top = Inches(0.1)
    set_first_bullet(tf_r, "Local AI (Like Our Lab)", size=Pt(21), bold=True, color=SECONDARY)
    for pt in [
        "Data stays on this machine \u2014 nowhere else.",
        "No internet connection needed to run.",
        "Full control over your data and models.",
        "You manage everything (updates, security).",
        "Privacy by design, but more responsibility.",
    ]:
        add_bullet(tf_r, pt, size=Pt(14), color=TEXT)

    callout = add_panel(slide, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.7),
                        fill=SURFACE, line=WARN)
    tfc = callout.text_frame
    tfc.margin_left = Inches(0.2)
    tfc.margin_top = Inches(0.06)
    set_first_bullet(
        tfc,
        "Always read the privacy policy. Some services opt you INTO training by default.",
        size=Pt(14), bold=True, color=WARN)

    add_footer(slide, "AI 101 Workshop")
    add_notes(
        slide,
        "SPEAKER NOTES — Privacy and Your Data\n"
        "• This is one of the most important slides — students need to hear this.\n"
        "• Ask: 'How many of you have used ChatGPT? Did you think about where your words go?'\n"
        "• Walk through both panels: cloud = your data travels, local = your data stays.\n"
        "• Emphasize the callout: some services USE your prompts for training by default.\n"
        "• Real scenario: 'If you paste your diary entry into ChatGPT, OpenAI employees might read it.'\n"
        "• Balance: 'Cloud AI isn't evil — but you should understand the tradeoff before using it.'\n"
        "• Transition: 'This connects directly to Lab 6 where you'll explore these tradeoffs in depth.'\n"
        "• Timing: ~3 min\n\n"
        "Image concept: shield split between cloud data flow and locked local vault.",
    )


def add_lab_intro_slide(prs, blank):
    """Slide: Welcome to the Lab!"""
    add_bullets_slide(
        prs, blank, "Welcome to the Lab!",
        [
            ("How to Connect", [
                "Open your browser \u2192 navigate to http://[server-ip]:3000",
                "No login required \u2014 just start chatting!"
            ]),
            ("Using the Interface", [
                "Model selector: pick different models from the dropdown.",
                "New chat: start fresh sessions for each exercise.",
                "Available models listed with descriptions in the sidebar."
            ]),
            ("What We'll Do", [
                "Hands-on exercises exploring different models and prompts.",
                "Compare outputs, speed, and quality across model sizes.",
                "Lab exercise guide available \u2014 follow along or explore!"
            ]),
        ],
        "students opening laptops and connecting to a glowing local server hub",
    )


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    add_title_slide(prs, blank)

    # 2
    add_bullets_slide(
        prs, blank, "Agenda",
        [
            "1. What AI is (and is not)",
            "2. How models are trained \u2014 from data to deployment",
            "3. How training and prediction really work",
            "4. Hallucinations \u2014 when AI makes stuff up",
            "5. Model types: text, multimodal, and image generation",
            "6. Model sizes, hardware, and choosing the right one",
            "7. Tokens, parameters, and context windows",
            "8. Agentic workflows, MCP, and tools",
            "9. Prompting patterns that work",
            "10. Hands-on lab \u2014 connect and experiment",
            "11. Safety, privacy, and responsible use",
            "12. Discussion and Q&A",
        ],
        "clean mission board agenda visual with neon accents",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Agenda\n"
        "• Walk through the agenda quickly — don't linger on each item.\n"
        "• Highlight: 'The best part is the hands-on lab where YOU get to talk to AI.'\n"
        "• Set pacing expectations: 'We have a lot to cover, but we'll keep it interactive.'\n"
        "• Ask: 'Has anyone here used an AI chatbot before? Raise your hand.'\n"
        "• Transition: 'Let's start with the big picture — what even IS AI?'\n"
        "• Timing: ~2 min")

    # 3
    add_bullets_slide(
        prs, blank, "AI in Layers: The Stack View",
        [
            ("Layer 1 - Models", ["The reasoning engine that predicts next tokens based on learned patterns."]),
            ("Layer 2 - Runtime", ["Local or cloud infrastructure that executes the model safely and fast."]),
            ("Layer 3 - Tools", ["File, browser, terminal, APIs, and external systems agents can call."]),
            ("Layer 4 - Product UX", ["Copilot/assistant experiences where users actually interact."]),
            ("Layer 5 - Governance", ["Policies, security controls, logs, and cost guardrails."]),
        ],
        "stacked architecture diagram with five labeled layers and arrows",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — AI in Layers\n"
        "• This slide gives the 30,000-foot view. Don't go deep — just set the mental model.\n"
        "• Walk bottom to top: 'Models are the brain, runtime is where they run, tools let them act.'\n"
        "• Relate to students: 'Today you interact at Layer 4 (the chat UI), but we'll peek under the hood.'\n"
        "• Transition: 'Let's zoom into Layer 1 — what IS a model?'\n"
        "• Timing: ~2 min")

    # 4
    add_bullets_slide(
        prs, blank, "What is a Model?",
        [
            "A model is a trained pattern engine, not a database of perfect facts.",
            "It maps input context to likely output tokens.",
            "Strengths: pattern recognition, synthesis, drafting, code transformation.",
            "Weaknesses: can hallucinate, can miss hidden constraints, no guaranteed truth.",
            "Better framing: probabilistic assistant, not oracle.",
        ],
        "glass brain made of text fragments and vectors, educational infographic style",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — What is a Model?\n"
        "• Key misconception to address: 'AI is NOT a search engine or database.'\n"
        "• Analogy: 'It's like someone who read the whole internet — they absorbed patterns, not facts.'\n"
        "• Ask: 'If I asked you to write a Shakespeare-sounding sentence, you could do it. That's what models do.'\n"
        "• Emphasize 'probabilistic assistant' — it gives you its best GUESS, not the truth.\n"
        "• Transition: 'So how do you actually build one of these pattern engines?'\n"
        "• Timing: ~2 min")

    # 5
    add_bullets_slide(
        prs, blank, "How Models are Built (High Level)",
        [
            ("Pretraining", ["Model learns language and code patterns from large datasets."]),
            ("Alignment / fine-tuning", ["Improves helpfulness, safety, and instruction following."]),
            ("Evaluation", ["Benchmarks + human feedback expose strengths and failure cases."]),
            ("Serving", ["Model is optimized, hosted, and monitored for real user workloads."]),
        ],
        "factory pipeline from data to trained model to deployment",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — How Models are Built\n"
        "• Keep it high-level — students don't need to understand backpropagation.\n"
        "• Factory analogy: 'Raw materials (data) → manufacturing (training) → quality check (eval) → shipping (serving).'\n"
        "• Key insight: 'The model learns patterns from examples — like learning a language by immersion.'\n"
        "• Transition: 'Let's go one level deeper into how training actually works.'\n"
        "• Timing: ~2 min")

    # 6 — How Training Really Works
    add_training_deep_dive_slide(prs, blank)
    override_last_slide_notes(prs,
        "SPEAKER NOTES — How Training Really Works\n"
        "• Three phases: data collection → pretraining → RLHF. Walk through each.\n"
        "• Data: 'Imagine downloading the entire internet — books, Wikipedia, code, forums.'\n"
        "• Pretraining: 'The model reads all of it and learns to predict the next word. That's basically it.'\n"
        "• RLHF: 'Humans grade the AI's answers: thumbs up or thumbs down. It learns from that feedback.'\n"
        "• Key insight: 'It learns PATTERNS, not FACTS. It might get things right because the pattern is common, not because it knows it's true.'\n"
        "• Transition: 'So how does prediction actually work, token by token?'\n"
        "• Timing: ~3 min")

    # 7 — How Does Prediction Actually Work?
    add_prediction_slide(prs, blank)
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Prediction\n"
        "• This is the 'aha moment' for most students.\n"
        "• Demo opportunity: 'I have a live token prediction demo — let me show you.'\n"
        "• Explain: 'The model picks the most likely next word, feeds it back, and repeats. Like autocomplete on steroids.'\n"
        "• Temperature: 'Low = boring but accurate. High = creative but risky. Like a spice dial.'\n"
        "• If running the token-prediction demo, pause and switch to it here.\n"
        "• Transition: 'What happens when this prediction engine gets it wrong?'\n"
        "• Timing: ~3 min (+ demo time)")

    # 8 — Hallucinations
    add_hallucinations_slide(prs, blank)
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Hallucinations\n"
        "• Start with: 'Raise your hand if you've ever seen AI confidently say something wrong.'\n"
        "• Key point: 'The model doesn't KNOW when it's wrong — it has no fact-checker inside.'\n"
        "• Share an example: fake book summaries, invented citations, wrong math.\n"
        "• Practical advice: 'Always fact-check important claims. Treat AI output like a first draft.'\n"
        "• Connect to Lab 4 where they'll catch hallucinations themselves.\n"
        "• Transition: 'Not all models just handle text — some can see images too.'\n"
        "• Timing: ~2 min")

    # 9 — Why Can Some Models Generate Images?
    add_multimodal_slide(prs, blank)
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Multimodal Models\n"
        "• Clarify the confusion: text models ≠ image generators. Different architectures.\n"
        "• Text model = trained on text. Multimodal = trained on text + images together.\n"
        "• Image generation (DALL-E, Stable Diffusion) = totally different process (diffusion).\n"
        "• Connect: 'In Lab 5, you'll use LLaVA — a model that can look at photos you upload.'\n"
        "• Fun fact: 'Some models can now understand video, audio, and even 3D objects.'\n"
        "• Transition: 'So where do these models actually run?'\n"
        "• Timing: ~2 min")

    # 10
    add_two_column_slide(
        prs, blank, "How Models Run: Local vs Cloud",
        "Local Model",
        [
            "Pros: privacy control, offline capability, predictable cost ceiling",
            "Pros: low-latency for small tasks when tuned well",
            "Cons: hardware limits, weaker large-model performance",
            "Cons: setup and maintenance burden",
        ],
        "Cloud Model",
        [
            "Pros: top-tier capability, elastic scaling, managed reliability",
            "Pros: easiest path to advanced multimodal and tool ecosystems",
            "Cons: usage-based cost and policy requirements",
            "Cons: network dependency",
        ],
        "split scene laptop edge device vs cloud datacenter pipeline",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Local vs Cloud\n"
        "• Point to the server: 'Everything today runs RIGHT HERE — nothing leaves this room.'\n"
        "• Walk through pros/cons of each side. Let students weigh in.\n"
        "• Ask: 'Would you rather have more privacy or more power? You'll explore this in Lab 6.'\n"
        "• Key insight: 'The gap is shrinking — local models are getting much better.'\n"
        "• Transition: 'Speaking of local models — size matters. Let's talk about that.'\n"
        "• Timing: ~3 min")

    # 7
    add_model_size_metaphor_slide(prs, blank)

    # 12 — What Hardware Runs a Model?
    add_hardware_slide(prs, blank)

    # 13
    add_two_column_slide(
        prs, blank, "Choose the Right Model for the Task",
        "Use Small/Fast Models For",
        [
            "Formatting, rewrite, summaries, classification",
            "Simple Q&A and low-risk automation",
            "High-volume tasks where latency matters",
            "Example: writing a haiku or polishing bullet points",
        ],
        "Use Larger Models For",
        [
            "Hard debugging and architecture design",
            "Cross-file code changes with nuanced constraints",
            "Multi-step reasoning with tradeoff analysis",
            "Example: designing robust auth refactors safely",
        ],
        "task-routing board with lightweight and heavyweight lanes",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Choose the Right Model\n"
        "• Key message: 'Match the tool to the job. Don't use a sledgehammer for a nail.'\n"
        "• Left side: small models are FAST and CHEAP. Perfect for simple tasks.\n"
        "• Right side: big models THINK DEEPER. Needed for complex work.\n"
        "• Ask: 'If you're writing a quick text, do you need the smartest AI? Probably not.'\n"
        "• Connect to Lab 2 where they'll compare models head-to-head.\n"
        "• Transition: 'Let's compare coding work vs other kinds of content.'\n"
        "• Timing: ~2 min")

    # 9
    add_two_column_slide(
        prs, blank, "Coding vs General Document Creation",
        "Coding Work",
        [
            "Needs correctness, tests, constraints, and execution checks",
            "Tool use matters: terminal, repo search, CI feedback",
            "Errors can ship to production if unchecked",
        ],
        "Docs / Content Work",
        [
            "Needs clarity, tone, and audience fit",
            "Fast iteration with human review is usually enough",
            "Lower technical blast radius, but still fact-check important claims",
        ],
        "split keyboard-and-code panel versus notebook-and-editor panel",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Coding vs Docs\n"
        "• Quick comparison: coding has strict correctness requirements. Docs are more forgiving.\n"
        "• For students who code: 'AI is amazing for debugging, but ALWAYS run the code to verify.'\n"
        "• For everyone: 'Even for writing, fact-check important claims.'\n"
        "• Transition: 'Now let's look at the tools available in the AI landscape.'\n"
        "• Timing: ~2 min")

    # 10
    add_tools_landscape_slide(prs, blank)

    # 11
    add_mcp_slide(prs, blank)

    # 12
    add_tokens_slide(prs, blank)

    # 18 — Parameters vs Tokens
    add_params_vs_tokens_slide(prs, blank)
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Parameters vs Tokens\n"
        "• Students often confuse these. This slide makes it crystal clear.\n"
        "• Parameters = the engine (fixed after training). Tokens = the fuel (flowing at runtime).\n"
        "• Car analogy: 'A V8 engine (big model) is more powerful, but uses more fuel (tokens cost more).'\n"
        "• Ask: 'If parameters are the engine size, what's the token limit?' (Answer: gas tank / context window)\n"
        "• Transition: 'Now let's talk about HOW to use these tools effectively — prompting.'\n"
        "• Timing: ~2 min")

    # 19
    add_bullets_slide(
        prs, blank, "Prompting Patterns That Work in Class",
        [
            ("Pattern 1 - Give role + goal + constraints", [
                "Example: 'Act as a code reviewer; find security issues only; keep it concise.'"
            ]),
            ("Pattern 2 - Ask for options with tradeoffs", [
                "Prompts students to think, not just accept first output."
            ]),
            ("Pattern 3 - Iterate with evidence", [
                "Ask model to cite files, lines, or assumptions before finalizing."
            ]),
            ("Pattern 4 - Keep a running summary", [
                "Prevents context drift in long sessions."
            ]),
        ],
        "teacher and students at whiteboard refining prompts together",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Prompting Patterns\n"
        "• This is the 'practical skills' slide. Students will USE these in the labs.\n"
        "• Pattern 1: 'Give the AI a job title and clear instructions. Like hiring a specialist.'\n"
        "• Pattern 2: 'Don't accept the first answer. Ask for alternatives and tradeoffs.'\n"
        "• Pattern 3: 'Make the AI show its work — ask it to cite sources or explain reasoning.'\n"
        "• Pattern 4: 'In long chats, periodically ask the AI to summarize what you've discussed.'\n"
        "• Connect directly to Lab 3 where they'll practice prompt engineering.\n"
        "• Transition: 'Time to put all of this into practice — welcome to the lab!'\n"
        "• Timing: ~3 min")

    # 20 — Welcome to the Lab!
    add_lab_intro_slide(prs, blank)
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Welcome to the Lab!\n"
        "• This is the transition to hands-on time. Energy should be high.\n"
        "• Share the server IP address now: write it on the board or project it.\n"
        "• Walk through: 'Open your browser, go to this URL, and you should see the chat interface.'\n"
        "• Troubleshoot: give students 2 minutes to connect. Help anyone having issues.\n"
        "• Point out the model dropdown and New Chat button.\n"
        "• Transition: 'Let's do a quick guided demo before you explore on your own.'\n"
        "• Timing: ~3 min")

    # 21
    add_bullets_slide(
        prs, blank, "Live Lab Plan (Student-Friendly)",
        [
            "Demo A: ask agent for a simple content rewrite (fast model candidate).",
            "Demo B: ask agent to propose a multi-file code fix (larger model candidate).",
            "Compare quality, speed, and cost assumptions as a class.",
            "Inspect what tools were used and whether MCP access was needed.",
            "Debrief: where humans must stay in control.",
        ],
        "classroom lab with students observing terminal and PR diff screens",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Live Lab Plan\n"
        "• Do Demo A live: pick a small model, show a simple rewrite. Fast and impressive.\n"
        "• Do Demo B live: pick a larger model, show a more complex task. Compare quality.\n"
        "• Ask the class: 'Which response was better? Was the extra wait time worth it?'\n"
        "• Then release students to work through the lab exercises at their own pace.\n"
        "• Circulate the room. Help students who get stuck. Celebrate cool discoveries.\n"
        "• Timing: ~5 min for demos, then student lab time")

    # 15
    add_bullets_slide(
        prs, blank, "Safety and Reality Check",
        [
            "AI can accelerate work, but does not replace engineering judgment.",
            "Always verify claims, code behavior, and external facts.",
            "Treat model output as draft material until reviewed.",
            "Use least-privilege tool access and clear repo policies.",
            "Responsible teams optimize for quality + safety, not just speed.",
        ],
        "seatbelt-on-rocket metaphor for safe acceleration",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Safety and Reality Check\n"
        "• Tone: serious but not scary. 'AI is powerful, and with power comes responsibility.'\n"
        "• Key message: 'Never blindly trust AI output. Always verify.'\n"
        "• Ask: 'What's the worst that could happen if AI gives wrong medical advice? Wrong code?'\n"
        "• Practical: 'Treat AI like a smart intern — great ideas, but needs supervision.'\n"
        "• Transition: 'Let's talk specifically about your data and privacy.'\n"
        "• Timing: ~2 min")

    # 23 — Privacy and Your Data
    add_privacy_slide(prs, blank)

    # 24
    add_storyboard_slide(prs, blank)

    # 17
    add_bullets_slide(
        prs, blank, "Discussion and Q&A",
        [
            "What felt most surprising?",
            "Where would you trust AI today, and where not yet?",
            "What workflow would you test first in this repo?",
            "What guardrails would make you comfortable using agents?",
        ],
        "open mic classroom discussion circle, collaborative energy",
    )
    override_last_slide_notes(prs,
        "SPEAKER NOTES — Discussion and Q&A\n"
        "• Open it up: 'What was the most surprising thing you learned today?'\n"
        "• Encourage honest answers. 'There are no wrong answers here.'\n"
        "• If the room is quiet, start with: 'I'll share mine — [your own takeaway].'\n"
        "• Good follow-ups: 'Would you use AI for homework? What about exams?'\n"
        "• Mention the feedback form if you set one up (QR code on screen).\n"
        "• Close with: 'AI is the most powerful tool you'll have in your careers. Learn it well.'\n"
        "• Timing: ~5-10 min")

    out_path = Path(__file__).parent / "Agentic_AI_101_Workshop.pptx"
    try:
        prs.save(str(out_path))
        print(f"Saved: {out_path}")
    except PermissionError:
        alt_path = Path(__file__).parent / "Agentic_AI_101_Workshop_updated.pptx"
        prs.save(str(alt_path))
        print(f"Saved (locked original, wrote alternate): {alt_path}")


if __name__ == "__main__":
    build_deck()
