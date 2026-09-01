#!/usr/bin/env python3
"""Generate the AI Foundations workshop slide deck.

Usage:
    python generate_workshopv2_deck.py
    python generate_workshopv2_deck.py --output AI_Workshop_V2.pptx

The content is intentionally curated for slides rather than copied paragraph-for-
paragraph from workshopv2.md. Edit the build_deck() calls to change slide wording
or order, and edit the theme constants to restyle the entire deck.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE_DOC = ROOT / "workshopv2.md"
DEFAULT_OUTPUT = ROOT / "AI_Workshop_V2.pptx"

W = Inches(13.333)
H = Inches(7.5)

# Theme: "Electric Blueprint"
INK = RGBColor(0x08, 0x10, 0x22)
INK_2 = RGBColor(0x0E, 0x1A, 0x33)
PANEL = RGBColor(0x14, 0x24, 0x42)
PANEL_2 = RGBColor(0x1B, 0x31, 0x55)
WHITE = RGBColor(0xF7, 0xFA, 0xFF)
TEXT = RGBColor(0xD7, 0xE2, 0xF2)
MUTED = RGBColor(0x8F, 0xA4, 0xC2)
CYAN = RGBColor(0x2D, 0xE2, 0xE6)
PURPLE = RGBColor(0xA7, 0x78, 0xFF)
CORAL = RGBColor(0xFF, 0x67, 0x76)
AMBER = RGBColor(0xFF, 0xC8, 0x57)
GREEN = RGBColor(0x55, 0xE6, 0x96)
BLUE = RGBColor(0x4C, 0x8D, 0xFF)

FONT_HEAD = "Bahnschrift SemiBold"
FONT_BODY = "Aptos"
FONT_MONO = "Cascadia Mono"


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def solid(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def set_line(shape, color: RGBColor, width: float = 1.5, transparency: int = 0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 24,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_text(
    slide,
    parts,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 24,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for text, color, bold in parts:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = PANEL,
    line: RGBColor | None = None,
    radius=True,
    transparency: int = 0,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    solid(shape, fill, transparency)
    if line:
        set_line(shape, line)
    return shape


def add_circle(
    slide,
    x: float,
    y: float,
    d: float,
    *,
    fill: RGBColor,
    transparency: int = 0,
    line: RGBColor | None = None,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    solid(shape, fill, transparency)
    if line:
        set_line(shape, line)
    return shape


def add_background(slide, accent: RGBColor = CYAN, variant: int = 0) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK

    # Restrained, single-accent geometry keeps the canvas modern without
    # competing with diagrams and text.
    add_box(slide, 11.75, -0.35, 1.95, 2.25, fill=accent, radius=True, transparency=80)
    add_box(slide, -0.45, 7.15, 2.4, 0.55, fill=accent, radius=True, transparency=88)
    add_box(slide, 0, 0, 0.1, 7.5, fill=accent, radius=False, transparency=8)
    divider = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.62),
        Inches(6.88),
        Inches(12.7),
        Inches(6.88),
    )
    set_line(divider, accent, 0.7, 80)


def add_kicker(slide, text: str, color: RGBColor, number: str | None = None) -> None:
    if number:
        add_box(slide, 0.55, 0.42, 0.5, 0.38, fill=color, radius=True)
        add_text(
            slide,
            number,
            0.55,
            0.42,
            0.5,
            0.38,
            size=11,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        x = 1.1
    else:
        x = 0.62
    add_text(slide, text.upper(), x, 0.47, 6.5, 0.28, size=11, color=color, bold=True)


def add_title(slide, title: str, subtitle: str | None = None, color: RGBColor = WHITE) -> None:
    add_text(slide, title, 0.62, 0.9, 12.0, 0.75, size=31, color=color, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.63, 11.7, 0.5, size=16, color=MUTED)


def add_footer(slide, index: int, label: str = "AI FOUNDATIONS") -> None:
    add_text(slide, label, 0.62, 7.11, 3.2, 0.2, size=8, color=MUTED, bold=True)
    add_text(slide, f"{index:02d}", 12.15, 7.08, 0.55, 0.23, size=9, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def new_slide(prs, accent=CYAN, variant=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent, variant)
    return slide


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    accent: RGBColor,
    *,
    badge: str | None = None,
    title_size: float = 18,
    body_size: float = 13,
    detail: str | None = None,
):
    add_box(slide, x, y, w, h, fill=PANEL, line=accent)
    add_box(slide, x, y, 0.08, h, fill=accent, radius=False)
    if badge:
        add_box(slide, x + 0.28, y + 0.28, 0.52, 0.46, fill=accent, radius=True)
        add_text(
            slide,
            badge,
            x + 0.28,
            y + 0.28,
            0.52,
            0.46,
            size=12,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        title_x = x + 0.95
        title_w = w - 1.2
    else:
        title_x = x + 0.32
        title_w = w - 0.55
    add_text(slide, title, title_x, y + 0.27, title_w, 0.48, size=title_size, color=WHITE, bold=True, font=FONT_HEAD)
    body_height = h - 1.92 if detail else h - 1.12
    add_text(slide, body, x + 0.32, y + 0.92, w - 0.58, body_height, size=body_size, color=TEXT)
    if detail:
        add_text(slide, "EXAMPLE", x + 0.32, y + h - 0.86, w - 0.58, 0.2, size=8, color=accent, bold=True)
        add_text(
            slide,
            detail,
            x + 0.32,
            y + h - 0.61,
            w - 0.58,
            0.38,
            size=9,
            color=WHITE,
            bold=True,
            font=FONT_MONO,
        )


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color=CYAN, width=2.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(line, color, width)
    line.line.end_arrowhead = True
    return line


def title_slide(prs, title, subtitle, notes):
    slide = new_slide(prs, CYAN)
    add_text(slide, "AI", 0.62, 0.38, 2.2, 1.2, size=84, color=CYAN, bold=True, font=FONT_HEAD)
    add_text(slide, "FOUNDATIONS", 2.65, 0.7, 5.3, 0.7, size=34, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(slide, title, 0.72, 2.15, 11.8, 1.5, size=48, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(slide, subtitle, 0.77, 3.85, 10.8, 0.7, size=22, color=TEXT)

    for i, (label, color) in enumerate([
        ("MODEL", CYAN),
        ("CONTEXT", PURPLE),
        ("TOOLS", CORAL),
    ]):
        y = 4.35 + i * 0.72
        add_box(slide, 9.05 + i * 0.35, y, 2.6, 0.52, fill=PANEL, line=color)
        add_box(slide, 9.05 + i * 0.35, y, 0.08, 0.52, fill=color, radius=False)
        add_text(slide, label, 9.35 + i * 0.35, y + 0.1, 2.0, 0.3, size=12, color=color, bold=True)

    add_text(slide, "For FIRST Robotics students | 2026", 0.77, 6.72, 5.0, 0.3, size=11, color=MUTED, bold=True)
    add_notes(slide, notes)


def act_slide(prs, act, title, subtitle, accent, notes):
    slide = new_slide(prs, accent, int(act))
    add_text(slide, f"ACT {act}", 0.7, 0.65, 2.2, 0.4, size=14, color=accent, bold=True)
    add_text(slide, title, 0.7, 1.35, 11.8, 1.3, size=54, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(slide, subtitle, 0.74, 2.95, 10.9, 0.8, size=23, color=TEXT)
    add_text(slide, f"0{act}", 9.45, 4.15, 2.75, 1.65, size=92, color=accent, bold=True, font=FONT_HEAD, align=PP_ALIGN.RIGHT)
    for i, width in enumerate([2.9, 2.2, 1.5]):
        add_box(slide, 8.95, 5.88 + i * 0.2, width, 0.06, fill=accent, radius=False, transparency=i * 18)
    add_notes(slide, notes)


def statement_slide(prs, index, kicker, statement, support, accent, notes):
    slide = new_slide(prs, accent, index)
    add_kicker(slide, kicker, accent, str(index))
    add_text(slide, statement, 0.72, 1.3, 11.8, 2.15, size=42, color=WHITE, bold=True, font=FONT_HEAD)
    add_box(slide, 0.72, 4.22, 11.9, 1.55, fill=PANEL, line=accent)
    add_text(slide, support, 1.05, 4.57, 11.2, 0.9, size=20, color=TEXT, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, index)
    add_notes(slide, notes)


def weights_definition_slide(prs, index, notes):
    slide = new_slide(prs, CYAN, index)
    add_kicker(slide, "Weights + parameters", CYAN, str(index))
    add_title(slide, "Weights are the model's learned settings")
    add_rich_text(
        slide,
        [
            ("WEIGHT", CYAN, True),
            (" = a parameter controlling influence     ", TEXT, False),
            ("PARAMETER", PURPLE, True),
            (" = any learned number     ", TEXT, False),
            ("TOKEN", AMBER, True),
            (" = a numbered piece of text", TEXT, False),
        ],
        0.75,
        1.72,
        11.75,
        0.35,
        size=13,
    )

    cards = [
        (
            "ARCHITECTURE",
            "FIXED MACHINERY",
            "Defines which mathematical operations happen and how signals move through the model.",
            "Think: engine design",
            CYAN,
        ),
        (
            "PARAMETERS",
            "LEARNED SETTINGS",
            "Training adjusts billions of numbers until useful language patterns emerge.",
            "Most are weights",
            PURPLE,
        ),
        (
            "TOKEN SCORES",
            "RUNTIME RESULT",
            "The machinery uses those learned settings to score possible next tokens.",
            'Think: "Paris" 95%',
            AMBER,
        ),
    ]
    for i, (label, heading, body, analogy, color) in enumerate(cards):
        x = 0.68 + i * 4.2
        add_box(slide, x, 2.45, 3.58, 3.55, fill=PANEL, line=color)
        add_text(slide, label, x + 0.3, 2.78, 2.8, 0.28, size=10, color=color, bold=True)
        add_text(slide, heading, x + 0.3, 3.24, 2.98, 0.42, size=18, color=WHITE, bold=True)
        add_text(slide, body, x + 0.3, 3.92, 2.98, 1.0, size=15, color=TEXT)
        add_box(slide, x + 0.3, 5.24, 2.98, 0.48, fill=INK_2, line=color)
        add_text(
            slide,
            analogy,
            x + 0.4,
            5.34,
            2.78,
            0.25,
            size=12,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if i < len(cards) - 1:
            add_arrow(slide, x + 3.68, 4.18, x + 4.08, 4.18, MUTED, 1.5)

    add_text(
        slide,
        "The architecture is the algorithm. The parameters are what training learns.",
        1.4,
        6.33,
        10.5,
        0.38,
        size=17,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, index)
    add_notes(slide, notes)


def weights_artifact_slide(prs, index, notes):
    slide = new_slide(prs, CYAN, index)
    add_kicker(slide, "Model package", CYAN, str(index))
    add_title(slide, "What one local model package looks like")
    add_text(
        slide,
        "One particular recipe, data mix, and training process produced this artifact.",
        0.75,
        1.72,
        11.4,
        0.35,
        size=15,
        color=TEXT,
    )

    add_box(slide, 0.65, 2.3, 5.15, 3.75, fill=PANEL, line=CYAN)
    add_text(slide, "MODEL PACKAGE SNAPSHOT", 1.0, 2.65, 3.5, 0.3, size=11, color=CYAN, bold=True)
    add_text(
        slide,
        "qwen3.8:27b-q4_K_M",
        1.0,
        3.17,
        4.45,
        0.55,
        size=22,
        color=WHITE,
        bold=True,
        font=FONT_MONO,
    )
    add_text(slide, "17 GB", 1.0, 4.0, 2.4, 0.75, size=42, color=CYAN, bold=True, font=FONT_HEAD)
    add_text(slide, "quantized model file", 3.0, 4.29, 2.25, 0.4, size=14, color=MUTED)
    add_text(
        slide,
        "~27 billion learned parameters",
        1.0,
        5.15,
        4.35,
        0.5,
        size=17,
        color=TEXT,
        bold=True,
    )

    details = [
        ("27b", "Parameter count", "About 27 billion learned values.", PURPLE),
        ("q4_K_M", "Quantization recipe", "Many weights stored with fewer bits.", AMBER),
        ("17 GB", "Package size", "The compressed file stored on disk.", CORAL),
    ]
    for i, (value, label, body, color) in enumerate(details):
        y = 2.3 + i * 1.25
        add_box(slide, 6.15, y, 6.5, 1.02, fill=PANEL, line=color)
        add_text(slide, value, 6.48, y + 0.22, 1.6, 0.42, size=19, color=color, bold=True, font=FONT_MONO)
        add_text(slide, label, 8.15, y + 0.18, 2.35, 0.3, size=14, color=WHITE, bold=True)
        add_text(slide, body, 8.15, y + 0.53, 4.05, 0.28, size=12, color=TEXT)

    add_box(slide, 6.15, 6.02, 6.5, 0.56, fill=INK_2, line=CYAN)
    add_text(
        slide,
        "A package of math values - not 27 billion readable facts.",
        6.42,
        6.13,
        5.95,
        0.26,
        size=13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, index)
    add_notes(slide, notes)


def database_exhibit_slide(prs, index, notes):
    slide = new_slide(prs, CYAN, index)
    add_kicker(slide, "Database", CYAN, str(index))
    add_title(slide, "A database stores values you can point to")

    x = 0.75
    y = 2.2
    widths = [1.45, 2.1, 0.75, 1.85, 1.5]
    headers = ["customerId", "name", "age", "location", "country"]
    rows = [
        ["CUST-001", "Ava Chen", "17", "Seattle", "USA"],
        ["CUST-002", "Kofi Mensah", "17", "Accra", "Ghana"],
        ["CUST-003", "Sofia Rossi", "16", "Milan", "Italy"],
    ]

    cursor = x
    for header, width in zip(headers, widths):
        add_box(slide, cursor, y, width, 0.72, fill=CYAN, radius=False)
        add_text(
            slide,
            header,
            cursor,
            y,
            width,
            0.72,
            size=13,
            color=INK,
            bold=True,
            font=FONT_MONO,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        cursor += width + 0.04

    for row_index, row in enumerate(rows):
        cursor = x
        row_y = y + 0.77 + row_index * 0.78
        for value, width in zip(row, widths):
            color = PANEL_2 if row_index == 1 else PANEL
            line = AMBER if row_index == 1 else MUTED
            add_box(slide, cursor, row_y, width, 0.72, fill=color, line=line, radius=False)
            add_text(
                slide,
                value,
                cursor,
                row_y,
                width,
                0.72,
                size=13,
                color=WHITE,
                bold=row_index == 1,
                font=FONT_MONO,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
            cursor += width + 0.04

    add_box(slide, 8.85, 2.2, 3.65, 3.06, fill=PANEL, line=AMBER)
    add_text(slide, "LOOKUP", 9.18, 2.52, 1.25, 0.3, size=11, color=AMBER, bold=True)
    add_text(
        slide,
        "customerId\n=\nCUST-002",
        9.18,
        3.0,
        2.98,
        1.35,
        size=21,
        color=WHITE,
        bold=True,
        font=FONT_MONO,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "returns the stored row", 9.18, 4.55, 2.98, 0.35, size=14, color=TEXT, align=PP_ALIGN.CENTER)

    add_box(slide, 1.55, 5.75, 10.2, 0.62, fill=INK_2, line=CYAN)
    add_text(
        slide,
        "The value is retrieved from a defined location - not predicted.",
        1.8,
        5.9,
        9.7,
        0.32,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, index)
    add_notes(slide, notes)


def vocabulary_scoring_slide(prs, index, notes):
    slide = new_slide(prs, CYAN, index)
    add_kicker(slide, "Vocabulary + scores", CYAN, str(index))
    add_title(slide, "How numbers become coherent text")

    add_box(slide, 0.72, 1.95, 4.0, 4.45, fill=PANEL, line=PURPLE)
    add_text(slide, "FIXED TOKEN VOCABULARY", 1.05, 2.24, 2.8, 0.3, size=11, color=PURPLE, bold=True)
    add_text(slide, "TOKEN ID", 1.08, 2.68, 0.9, 0.24, size=9, color=MUTED, bold=True)
    add_text(slide, "TEXT PIECE", 2.48, 2.68, 1.3, 0.24, size=9, color=MUTED, bold=True)
    vocab = [
        ("#91", '"what"'),
        ("#522", '"is"'),
        ("#2210", '"capital"'),
        ("#1188", '"France"'),
        ("#4402", '"Paris"'),
    ]
    for i, (token_id, token_text) in enumerate(vocab):
        row_y = 3.02 + i * 0.53
        add_text(slide, token_id, 1.08, row_y, 0.9, 0.32, size=15, color=CYAN, bold=True, font=FONT_MONO)
        add_text(slide, "->", 2.0, row_y, 0.42, 0.32, size=15, color=MUTED, bold=True, font=FONT_MONO)
        add_text(slide, token_text, 2.48, row_y, 1.75, 0.32, size=15, color=WHITE, bold=True, font=FONT_MONO)
    add_text(
        slide,
        "IDs are labels - not weights or parameters.",
        1.08,
        5.72,
        3.25,
        0.42,
        size=13,
        color=TEXT,
        bold=True,
    )

    add_box(slide, 4.98, 1.95, 7.65, 4.45, fill=PANEL, line=CYAN)
    add_text(slide, "QUESTION IN", 5.32, 2.27, 1.8, 0.3, size=11, color=CYAN, bold=True)
    add_text(
        slide,
        '"what is the capital of France?"',
        5.32,
        2.69,
        6.7,
        0.42,
        size=18,
        color=WHITE,
        bold=True,
        font=FONT_MONO,
    )
    add_rich_text(
        slide,
        [
            ("ARCHITECTURE", MUTED, True),
            (" + ", TEXT, False),
            ("LEARNED PARAMETERS", PURPLE, True),
            (" -> NEXT-TOKEN SCORES", AMBER, True),
        ],
        5.32,
        3.37,
        6.8,
        0.3,
        size=11,
    )

    scores = [
        ("Paris", "95.0%", 5.1, CYAN),
        ("London", "2.1%", 0.22, PURPLE),
        ("Lyon", "1.3%", 0.15, CORAL),
        ("Nice", "0.4%", 0.1, AMBER),
    ]
    for i, (label, value, bar_width, color) in enumerate(scores):
        row_y = 3.82 + i * 0.52
        add_text(slide, label, 5.32, row_y, 1.0, 0.28, size=13, color=WHITE, bold=True, font=FONT_MONO)
        add_box(slide, 6.38, row_y + 0.02, bar_width, 0.25, fill=color, radius=True)
        add_text(slide, value, 11.65, row_y, 0.55, 0.28, size=12, color=color, bold=True, font=FONT_MONO, align=PP_ALIGN.RIGHT)

    add_text(slide, 'CHOSEN TOKEN:  "Paris"', 5.32, 5.98, 3.25, 0.35, size=17, color=CYAN, bold=True, font=FONT_MONO)
    add_text(
        slide,
        "ILLUSTRATIVE EXAMPLE - TOKEN IDs AND SCORES ARE PLACEHOLDERS",
        4.95,
        6.5,
        7.7,
        0.22,
        size=8,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_footer(slide, index)
    add_notes(slide, notes)


def three_cards_slide(prs, index, kicker, title, cards, accent, notes):
    slide = new_slide(prs, accent, index)
    add_kicker(slide, kicker, accent, str(index))
    add_title(slide, title)
    for i, (head, body, color) in enumerate(cards):
        add_card(slide, 0.62 + i * 4.18, 2.35, 3.82, 3.72, head, body, color, badge=str(i + 1))
    add_footer(slide, index)
    add_notes(slide, notes)


def compare_slide(prs, index, kicker, title, left, right, accent, notes):
    slide = new_slide(prs, accent, index)
    add_kicker(slide, kicker, accent, str(index))
    add_title(slide, title)
    for x, data, color in [(0.62, left, CYAN), (6.88, right, PURPLE)]:
        heading, tag, points = data
        add_box(slide, x, 2.16, 5.83, 4.25, fill=PANEL, line=color)
        add_text(slide, tag.upper(), x + 0.35, 2.43, 2.4, 0.3, size=10, color=color, bold=True)
        add_text(slide, heading, x + 0.35, 2.83, 5.1, 0.6, size=27, color=WHITE, bold=True, font=FONT_HEAD)
        for j, point in enumerate(points):
            add_circle(slide, x + 0.38, 3.75 + j * 0.72, 0.18, fill=color)
            add_text(slide, point, x + 0.72, 3.61 + j * 0.72, 4.72, 0.52, size=15, color=TEXT)
    add_footer(slide, index)
    add_notes(slide, notes)


def flow_slide(prs, index, kicker, title, steps, accent, notes, bottom=None, label=None):
    slide = new_slide(prs, accent, index)
    add_kicker(slide, kicker, accent, str(index))
    add_title(slide, title)
    if label:
        add_text(slide, label, 0.75, 2.08, 5.4, 0.28, size=10, color=accent, bold=True)
    count = len(steps)
    gap = 0.23
    width = (12.05 - gap * (count - 1)) / count
    y = 2.55
    for i, step in enumerate(steps):
        if len(step) == 4:
            head, body, detail, color = step
        else:
            head, body, color = step
            detail = None
        x = 0.64 + i * (width + gap)
        add_card(
            slide,
            x,
            y,
            width,
            2.75,
            head,
            body,
            color,
            badge=str(i + 1),
            title_size=16,
            body_size=12,
            detail=detail,
        )
        if i < count - 1:
            add_arrow(slide, x + width + 0.02, y + 1.37, x + width + gap - 0.04, y + 1.37, MUTED, 1.5)
    if bottom:
        add_box(slide, 1.15, 5.72, 11.05, 0.62, fill=INK_2, line=accent)
        add_text(slide, bottom, 1.35, 5.85, 10.65, 0.35, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def bubbles_slide(prs, index, kicker, title, bubbles, accent, notes, center=None):
    slide = new_slide(prs, accent, index)
    add_kicker(slide, kicker, accent, str(index))
    add_title(slide, title)
    if center:
        add_circle(slide, 5.35, 2.66, 2.55, fill=PANEL_2, line=accent)
        add_text(slide, center, 5.53, 3.32, 2.18, 0.85, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    positions = [(1.05, 2.5), (9.45, 2.35), (1.85, 4.85), (8.65, 4.85)]
    for (label, body, color), (x, y) in zip(bubbles, positions):
        add_circle(slide, x, y, 2.05, fill=color, transparency=8)
        add_text(slide, label, x + 0.18, y + 0.34, 1.69, 0.38, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.2, y + 0.8, 1.65, 0.88, size=10.5, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if center:
            add_arrow(slide, x + 1.02, y + 1.02, 6.62, 3.92, color, 1.3)
    add_footer(slide, index)
    add_notes(slide, notes)


def stack_slide(prs, index, kicker, title, layers, accent, notes):
    slide = new_slide(prs, accent, index)
    add_kicker(slide, kicker, accent, str(index))
    add_title(slide, title)
    for i, (name, description, color) in enumerate(layers):
        y = 2.25 + i * 0.82
        x = 1.0 + i * 0.16
        width = 11.25 - i * 0.32
        add_box(slide, x, y, width, 0.66, fill=PANEL, line=color)
        add_text(slide, name, x + 0.28, y + 0.12, 2.5, 0.34, size=15, color=color, bold=True)
        add_text(slide, description, x + 2.85, y + 0.12, width - 3.12, 0.34, size=13, color=TEXT)
    add_footer(slide, index)
    add_notes(slide, notes)


def token_slide(prs, index, notes):
    slide = new_slide(prs, PURPLE, index)
    add_kicker(slide, "Tokens", PURPLE, str(index))
    add_title(slide, "The model does not see words the way you do")
    sentence = ["The", "driver", "said", "'", "let", "'s", "go", "'", "!"]
    colors = [CYAN, PURPLE, BLUE, CORAL, AMBER, GREEN, CYAN, CORAL, PURPLE]
    x = 0.7
    for token, color in zip(sentence, colors):
        width = max(0.65, 0.24 * len(token) + 0.42)
        add_box(slide, x, 2.55, width, 0.82, fill=color, radius=True)
        add_text(slide, token, x, 2.55, width, 0.82, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        x += width + 0.16
    add_rich_text(
        slide,
        [
            ("Rough English estimate: ", TEXT, False),
            ("1 token", PURPLE, True),
            (" is about ", TEXT, False),
            ("4 characters", CYAN, True),
            (" or ", TEXT, False),
            ("0.75 words", AMBER, True),
        ],
        1.0,
        4.25,
        11.4,
        0.62,
        size=23,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "Limits, speed, and context are measured in tokens - not words.", 1.2, 5.25, 10.9, 0.55, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def context_window_slide(prs, index, notes):
    slide = new_slide(prs, PURPLE, index)
    add_kicker(slide, "Context", PURPLE, str(index))
    add_title(slide, "A finite desk, not permanent memory")
    add_box(slide, 0.85, 2.15, 11.6, 3.65, fill=PANEL, line=PURPLE)
    add_text(slide, "CONTEXT WINDOW", 1.15, 2.42, 3.1, 0.35, size=12, color=PURPLE, bold=True)
    labels = [
        ("SYSTEM", 1.18, 3.08, 1.65, CYAN),
        ("HISTORY", 3.0, 3.08, 2.2, BLUE),
        ("FILES", 5.37, 3.08, 1.65, GREEN),
        ("TOOLS", 7.19, 3.08, 1.65, CORAL),
        ("YOUR ASK", 9.01, 3.08, 2.15, AMBER),
    ]
    for label, x, y, width, color in labels:
        add_box(slide, x, y, width, 1.15, fill=color)
        add_text(slide, label, x, y, width, 1.15, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "When the desk fills: compress, retrieve selectively, or drop older detail.", 1.15, 4.75, 10.9, 0.52, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def bob_slide(prs, index, notes):
    slide = new_slide(prs, PURPLE, index)
    add_kicker(slide, "Context", PURPLE, str(index))
    add_title(slide, "Why did it forget Bob?")
    stages = [
        ("1", "You say", '"My name is Bob."', CYAN),
        ("2", "Hours of coding", "Errors, fixes, files, decisions...", BLUE),
        ("3", "Compression", "Keep the project. Lose the small detail.", PURPLE),
        ("4", "You ask again", '"What is my name?"', AMBER),
        ("5", "No evidence", "Bob is no longer in the supplied context.", CORAL),
    ]
    for i, (num, head, body, color) in enumerate(stages):
        x = 0.6 + i * 2.52
        add_circle(slide, x + 0.72, 2.55, 0.68, fill=color)
        add_text(slide, num, x + 0.72, 2.55, 0.68, 0.68, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, head, x, 3.42, 2.14, 0.42, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x, 3.98, 2.14, 1.25, size=12, color=TEXT, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_arrow(slide, x + 1.63, 2.9, x + 2.42, 2.9, MUTED, 1.4)
    add_box(slide, 2.05, 5.73, 9.25, 0.58, fill=INK_2, line=PURPLE)
    add_text(slide, '"Remembering" means the system supplied the information again.', 2.2, 5.86, 8.95, 0.3, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def current_information_slide(prs, index, notes):
    slide = new_slide(prs, PURPLE, index)
    add_kicker(slide, "Current information", PURPLE, str(index))
    add_title(slide, "Stale at the core. Fresh at the edges.")

    stages = [
        (
            "MODEL CORE",
            "Uses the same frozen weights",
            "It reads new evidence without being retrained.",
            CYAN,
        ),
        (
            "RUNTIME CONTEXT",
            "Returned evidence is supplied now",
            "The result becomes input for the next model call.",
            PURPLE,
        ),
        (
            "HARNESS",
            "Runs a tool or calls an MCP server",
            "The surrounding application performs the retrieval.",
            CORAL,
        ),
        (
            "LIVE SOURCE",
            "Web, API, database, or files",
            "May contain information created after training.",
            GREEN,
        ),
    ]
    for i, (label, heading, body, color) in enumerate(stages):
        x = 0.55 + i * 3.05
        add_box(slide, x, 2.25, 2.7, 3.05, fill=PANEL, line=color)
        add_text(slide, label, x + 0.25, 2.55, 2.2, 0.28, size=10, color=color, bold=True)
        add_text(slide, heading, x + 0.25, 3.02, 2.2, 0.72, size=17, color=WHITE, bold=True)
        add_text(slide, body, x + 0.25, 4.02, 2.2, 0.82, size=13, color=TEXT)
        if i < len(stages) - 1:
            next_color = stages[i + 1][3]
            add_text(
                slide,
                "<-",
                x + 2.72,
                3.62,
                0.38,
                0.32,
                size=17,
                color=next_color,
                bold=True,
                font=FONT_MONO,
                align=PP_ALIGN.CENTER,
            )

    add_box(slide, 1.0, 5.75, 11.35, 0.66, fill=INK_2, line=PURPLE)
    add_rich_text(
        slide,
        [
            ("CURRENT ANSWER", WHITE, True),
            (" = ", TEXT, False),
            ("FROZEN WEIGHTS", CYAN, True),
            (" + ", TEXT, False),
            ("FRESH SUPPLIED EVIDENCE", GREEN, True),
        ],
        1.25,
        5.92,
        10.85,
        0.3,
        size=15,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, index)
    add_notes(slide, notes)


def hallucination_slide(prs, index, notes):
    slide = new_slide(prs, PURPLE, index)
    add_kicker(slide, "Hallucination", PURPLE, str(index))
    add_title(slide, "The model does not automatically fact-check itself")
    add_text(
        slide,
        "A hallucination is generated content that is false or unsupported.",
        0.75,
        1.72,
        11.7,
        0.42,
        size=17,
        color=TEXT,
    )

    cards = [
        (
            "SAME PROCESS",
            "Right or wrong, every answer is built one predicted token at a time.",
            CYAN,
        ),
        (
            "CHECKS VARY",
            "Some products add search, citations, tests, or review. Others show the generated answer directly.",
            PURPLE,
        ),
        (
            "YOUR RULE",
            "For anything important, confirm it with a trusted source, test, or repeatable calculation.",
            CORAL,
        ),
    ]
    for i, (heading, body, color) in enumerate(cards):
        add_card(
            slide,
            0.62 + i * 4.18,
            2.42,
            3.82,
            3.12,
            heading,
            body,
            color,
            badge=str(i + 1),
            title_size=16,
            body_size=14,
        )

    add_box(slide, 1.05, 5.86, 11.2, 0.62, fill=INK_2, line=AMBER)
    add_text(
        slide,
        "No separate check occurred? Treat the answer as unverified.",
        1.3,
        6.0,
        10.7,
        0.32,
        size=15,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, index)
    add_notes(slide, notes)


def tool_call_slide(prs, index, notes):
    slide = new_slide(prs, CORAL, index)
    add_kicker(slide, "Tools", CORAL, str(index))
    add_title(slide, "A tool call, end to end")
    nodes = [
        ("MODEL", 'web_search("FRC 2026")', CYAN),
        ("HARNESS", "Checks permission\nand runs it", PURPLE),
        ("TOOL", "Returns current\nsearch results", CORAL),
        ("MODEL", "Reads result\nand answers", GREEN),
    ]
    for i, (head, body, color) in enumerate(nodes):
        x = 0.62 + i * 3.18
        add_circle(slide, x + 0.72, 2.42, 1.4, fill=color)
        add_text(slide, head, x + 0.72, 2.42, 1.4, 1.4, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x, 4.05, 2.85, 1.08, size=14, color=TEXT, align=PP_ALIGN.CENTER, font=FONT_MONO if i == 0 else FONT_BODY)
        if i < len(nodes) - 1:
            add_arrow(slide, x + 2.2, 3.12, x + 3.02, 3.12, color, 2)
    add_box(slide, 1.15, 5.63, 11.05, 0.72, fill=PANEL, line=CORAL)
    add_text(slide, "The model proposes. The harness permits and acts.", 1.35, 5.82, 10.65, 0.34, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def mcp_slide(prs, index, notes):
    slide = new_slide(prs, CORAL, index)
    add_kicker(slide, "MCP", CORAL, str(index))
    add_title(slide, "One shared plug for many capabilities")
    add_circle(slide, 5.2, 2.33, 2.8, fill=PANEL_2, line=CORAL)
    add_text(slide, "AI HOST", 5.2, 2.72, 2.8, 0.5, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "+ MCP client", 5.2, 3.33, 2.8, 0.4, size=13, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("GITHUB", 0.8, 2.2, PURPLE),
        ("FILES", 1.55, 4.9, CYAN),
        ("DATABASE", 9.75, 2.2, GREEN),
        ("BROWSER", 9.05, 4.9, AMBER),
    ]
    for label, x, y, color in items:
        add_box(slide, x, y, 2.75, 1.02, fill=color)
        add_text(slide, label, x, y, 2.75, 1.02, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_arrow(slide, x + 1.37, y + 0.51, 6.6, 3.73, color, 1.8)
    add_text(slide, "Servers expose tools, resources, and prompts. The host chooses what the model may use.", 2.0, 6.28, 9.35, 0.5, size=16, color=TEXT, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def failure_slide(prs, index, notes):
    slide = new_slide(prs, AMBER, index)
    add_kicker(slide, "Failure modes", AMBER, str(index))
    add_title(slide, "When things go off the rails")
    failures = [
        ("LOOP", "It never decides to stop", "Add iteration and time limits", CORAL),
        ("INJECTION", "Untrusted data acts like an order", "Restrict, validate, approve", AMBER),
        ("HALLUCINATION", "Plausible becomes false", "Ground, test, verify", PURPLE),
        ("OVERLOAD", "Important context disappears", "Compress, retrieve, restate", BLUE),
    ]
    positions = [(0.7, 2.35), (6.82, 2.35), (0.7, 4.45), (6.82, 4.45)]
    for (name, symptom, fix, color), (x, y) in zip(failures, positions):
        add_box(slide, x, y, 5.8, 1.65, fill=PANEL, line=color)
        add_text(slide, name, x + 0.28, y + 0.23, 1.8, 0.3, size=11, color=color, bold=True)
        add_text(slide, symptom, x + 0.28, y + 0.62, 3.25, 0.6, size=17, color=WHITE, bold=True)
        add_text(slide, fix, x + 3.62, y + 0.54, 1.85, 0.72, size=12, color=TEXT, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, index)
    add_notes(slide, notes)


def recap_slide(prs, index, notes):
    slide = new_slide(prs, CYAN, index)
    add_kicker(slide, "Takeaway", CYAN, str(index))
    add_title(slide, "The whole system in one picture")
    pieces = [
        ("MODEL", "Predicts the next token", CYAN),
        ("CONTEXT", "Supplies the working information", PURPLE),
        ("HARNESS", "Builds the loop and enforces rules", BLUE),
        ("TOOLS", "Connect the loop to the world", CORAL),
        ("HUMAN", "Chooses goals and verifies results", AMBER),
    ]
    x = 0.55
    for i, (head, body, color) in enumerate(pieces):
        add_box(slide, x, 2.55, 2.28, 2.62, fill=PANEL, line=color)
        add_circle(slide, x + 0.76, 2.9, 0.76, fill=color)
        add_text(slide, str(i + 1), x + 0.76, 2.9, 0.76, 0.76, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, head, x + 0.2, 3.87, 1.88, 0.35, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.2, 4.32, 1.88, 0.62, size=11, color=TEXT, align=PP_ALIGN.CENTER)
        if i < len(pieces) - 1:
            add_arrow(slide, x + 2.29, 3.86, x + 2.56, 3.86, MUTED, 1.3)
        x += 2.58
    add_text(slide, "The product you experience is all five - not the model alone.", 1.2, 5.85, 10.9, 0.55, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, index)
    add_notes(slide, notes)


def closing_slide(prs, notes):
    slide = new_slide(prs, CYAN, 1)
    add_text(slide, "Stay curious.", 0.7, 1.0, 8.7, 0.9, size=52, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(slide, "Stay skeptical.", 0.7, 2.05, 8.7, 0.9, size=52, color=PURPLE, bold=True, font=FONT_HEAD)
    add_text(slide, "Build boldly.", 0.7, 3.1, 8.7, 0.9, size=52, color=CORAL, bold=True, font=FONT_HEAD)
    add_box(slide, 8.75, 1.15, 3.75, 4.7, fill=PANEL, line=CYAN)
    add_text(slide, "BEFORE YOU TRUST IT", 9.15, 1.58, 2.95, 0.45, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    for i, text in enumerate(["Check the source", "Run the code", "Test the output", "Own the decision"]):
        add_circle(slide, 9.15, 2.38 + i * 0.78, 0.35, fill=[CYAN, GREEN, PURPLE, AMBER][i])
        add_text(slide, text, 9.68, 2.27 + i * 0.78, 2.3, 0.45, size=15, color=WHITE, bold=True)
    add_text(slide, "Questions?", 0.75, 5.55, 5.0, 0.75, size=35, color=TEXT, bold=True, font=FONT_HEAD)
    add_text(slide, "AI Foundations for FIRST Robotics students", 0.78, 6.65, 6.2, 0.3, size=11, color=MUTED)
    add_notes(slide, notes)


def build_deck(output: Path) -> Path:
    if not SOURCE_DOC.exists():
        raise FileNotFoundError(f"Source document not found: {SOURCE_DOC}")

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "AI Foundations: How Modern AI Systems Actually Work"
    prs.core_properties.subject = "AI and agentic tool foundations for FIRST Robotics students"
    prs.core_properties.author = "Matt Vevang"
    prs.core_properties.keywords = "AI, LLM, context, tools, MCP, agents, FIRST Robotics"

    title_slide(
        prs,
        "How modern AI systems actually work",
        "Models, context, tools, agents, and the failure modes between them",
        "Welcome the team. Frame this as a mental-model workshop, not a product tutorial.\n"
        "The products will change. These concepts will remain useful.\n"
        "Source: workshopv2.md",
    )

    index = 2
    flow_slide(
        prs,
        index,
        "The map",
        "Four acts. One connected system.",
        [
            ("THE MODEL", "Weights, training, tokens, and generation.", CYAN),
            ("CONTEXT", "Working memory, current information, and hallucinations.", PURPLE),
            ("AGENTIC", "Harnesses, tools, MCP, instructions, and skills.", CORAL),
            ("TRUST + CONTROL", "Failure modes, limits, verification, and human responsibility.", AMBER),
        ],
        CYAN,
        "PRESENTER LEGEND - THE FOUR ACTS\n\n"
        "These act names are an editorial teaching structure for this workshop, not a "
        "standard industry taxonomy. Each act depends on concepts introduced earlier.\n\n"
        "ACT 1 - THE MODEL (slides 5-18)\n"
        "What happens inside the language model: databases as a contrast, prediction, "
        "weights, vocabulary and scoring, temperature, training, quantization, tokens, "
        "inference, and generation.\n\n"
        "ACT 2 - CONTEXT + TRUTH (slides 19-25)\n"
        "What information the model receives at runtime, why it appears to remember, how "
        "context is lost, how current information reaches it, and why hallucinations require "
        "verification.\n\n"
        "ACT 3 - AGENTIC SYSTEMS (slides 26-36)\n"
        "What surrounds the model: the harness, repeated agent loop, tools, file boundaries, "
        "MCP, instruction layers, system prompts, and reusable skills.\n\n"
        "ACT 4 - TRUST + CONTROL (slides 37-41)\n"
        "This is an editorial umbrella for operating capable systems responsibly. It covers "
        "failure modes, prompt injection, permissions, stopping limits, verification, and "
        "human accountability. 'Trust + Control' is not a named protocol or product feature.\n\n"
        "Use this slide as your navigation legend throughout the presentation.",
        "The product you experience is the complete system - not the model alone.",
    )
    index += 1

    statement_slide(
        prs,
        index,
        "What is AI?",
        "Software that uses learned patterns to produce predictions, recommendations, decisions, or content.",
        "This workshop focuses on large language models and the products built around them.",
        CYAN,
        "Keep the definition broad, then narrow the scope to LLM-based chat and coding assistants.",
    )
    index += 1

    three_cards_slide(
        prs,
        index,
        "What is the product?",
        "The product is the whole system",
        [
            ("MODEL", "The pattern engine. Learned weights turn context into token scores.", CYAN),
            ("HARNESS", "The application loop. It assembles context, applies rules, and decides what runs.", PURPLE),
            ("TOOLS + SERVICES", "The outside capabilities: files, search, browsers, APIs, databases, and more.", CORAL),
        ],
        CYAN,
        "Define the product positively: the experience students use is the model, harness, "
        "and connected capabilities working together. Reuse these three labels throughout.",
    )
    index += 1

    act_slide(
        prs,
        1,
        "Inside the model",
        "How learned numbers become one predicted token at a time",
        CYAN,
        "Transition: begin with a deterministic system students already understand, then contrast it with an LLM.",
    )
    index += 1

    database_exhibit_slide(
        prs,
        index,
        "Walk left to right. The table has defined columns and stored values. The highlighted "
        "customerId is a unique key that points to one row. Asking for CUST-002 retrieves "
        "Kofi Mensah from that stored location.\n\n"
        "This is the deterministic-side anchor for the rest of Act 1: same stored data and "
        "same lookup produce the same stored result.",
    )
    index += 1

    statement_slide(
        prs,
        index,
        "The sentence to remember",
        "An LLM is a very fancy weighted pattern matcher.",
        "It was trained by imitation and answers by prediction, not lookup.",
        CYAN,
        "Say 'fancy word guesser' if it helps, then immediately upgrade the phrase to weighted pattern matcher.",
    )
    index += 1

    weights_definition_slide(
        prs,
        index,
        "Define the terms before moving on. A parameter is any numeric setting learned during "
        "training. A weight is the most common kind of parameter; it controls how strongly one "
        "internal signal influences another. People often use 'weights' and 'parameters' almost "
        "interchangeably because weights make up nearly all of the parameter count.\n\n"
        "A token is not a parameter. It is a numbered piece of text supplied to or produced by "
        "the model. Tokens change with each prompt and response; the trained parameters remain "
        "fixed during ordinary generation.\n\n"
        "Walk left to right. The architecture is the mathematical machinery or program "
        "structure. The parameters are the learned settings inside that machinery. Given "
        "input tokens, the architecture uses those parameters to produce next-token scores.\n\n"
        "This distinction prepares the vocabulary, temperature, and training slides that follow.",
    )
    index += 1

    vocabulary_scoring_slide(
        prs,
        index,
        "This is an illustrative teaching example. The token IDs and percentages are invented "
        "to make the mechanism visible; they are not output copied from a real tokenizer or model.\n\n"
        "A token is a piece of text; the later token slide expands that definition. The left "
        "panel is the fixed vocabulary mapping between token IDs and text pieces. The vocabulary "
        "says what pieces are available, but it does not contain the fact that Paris is the "
        "capital of France. The integers such as #91 and #4402 are token IDs - addresses in "
        "that lookup table. They are not weights, parameters, or IDs for parameters.\n\n"
        "Keep the roles separate: the architecture is the mathematical machinery. The learned "
        "parameters are billions of settings used by that machinery. Together they calculate "
        "a score for every possible next-token ID. The vocabulary translates the chosen ID "
        "back into a text piece.\n\n"
        "The right panel shows the model scoring possible next tokens for the question. The "
        "learned weights make Paris score much higher than London, Lyon, or Nice. The system "
        "chooses Paris, appends it to the conversation, and repeats the process for the next token.\n\n"
        "Temperature comes after scoring. It changes how strongly selection favors the highest "
        "scores; it is not where the learned language relationships are stored.",
    )
    index += 1

    compare_slide(
        prs,
        index,
        "Temperature",
        "Temperature controls how adventurous each choice is",
        (
            "LOW TEMPERATURE",
            "FOLLOW THE FAVORITES",
            [
                "Mostly chooses the strongest candidate",
                "More repeatable and focused",
                "Zero is still not a perfect guarantee",
            ],
        ),
        (
            "HIGH TEMPERATURE",
            "EXPLORE ALTERNATIVES",
            [
                "Gives weaker candidates more of a chance",
                "More varied or creative",
                "More opportunity for mistakes",
            ],
        ),
        CYAN,
        "Connect directly to the previous scoring slide. Paris had the strongest score, while "
        "London and Lyon were weaker candidates.\n\n"
        "LOW TEMPERATURE stays close to the strongest candidate, so Paris is chosen almost "
        "every time. HIGH TEMPERATURE gives lower-scoring alternatives more opportunity, "
        "which can add variety but can also take the answer down a worse path.\n\n"
        "Temperature does not add knowledge or rewrite the scores. It changes how strictly "
        "the system follows those scores while choosing. The range and availability vary by "
        "model and provider.",
    )
    index += 1

    flow_slide(
        prs,
        index,
        "Training",
        "Training is how the weights are learned",
        [
            ("EXAMPLE", "Give the model text with a known next token.", CYAN),
            ("PREDICT", "Ask what token should come next.", PURPLE),
            ("CHECK", "Measure the prediction against the known answer.", AMBER),
            ("ADJUST", "Nudge the weights toward a better prediction.", CORAL),
            ("REPEAT", "Do it across an enormous training mixture.", GREEN),
        ],
        CYAN,
        "TRAINING - PRESENTER EXPLANATION\n\n"
        "This happens before the released model is used for chat. The training example includes "
        "a known next token, so the system can measure how wrong the prediction was. That error "
        "is used to make tiny adjustments to the weights.\n\n"
        "Repeat that process across an enormous training mixture and the weights gradually encode "
        "useful statistical relationships.\n\n"
        "During normal generation there is no supplied correct answer to compare against, and the "
        "weights remain fixed. The next slide puts the two loops side by side.\n\n"
        "Avoid a backpropagation or calculus detour unless someone specifically asks.",
        "TRAINING CHANGES THE WEIGHTS. GENERATION DOES NOT.",
    )
    index += 1

    compare_slide(
        prs,
        index,
        "Training vs. generation",
        "Two loops. Only one changes the model.",
        (
            "TRAINING",
            "BEFORE RELEASE",
            [
                "Uses examples with a known target",
                "Compares prediction with the target",
                "Adjusts the weights after errors",
            ],
        ),
        (
            "GENERATION",
            "WHEN YOU USE IT",
            [
                "Uses your prompt and current context",
                "Chooses and appends one token",
                "The trained weights stay fixed",
            ],
        ),
        CYAN,
        "The two loops look related because both ask the model to predict tokens.\n\n"
        "TRAINING has an answer key. It knows the target token, measures the error, and updates "
        "the weights. Its output is the trained model package.\n\n"
        "GENERATION has no answer key. It uses the fixed weights to score possible next tokens, "
        "chooses one, appends it, and repeats. Its output is the response the user sees.\n\n"
        "Analogy: the architecture is closer to the engine design or program structure. The "
        "weights are billions of learned settings inside that engine. A distributed model package "
        "bundles the architecture information, tokenizer, and learned weights together.",
    )
    index += 1

    three_cards_slide(
        prs,
        index,
        "Different models",
        "Similar ingredients. Different kitchens.",
        [
            ("THE RECIPE", "Architecture and scale: how the math is arranged and how much capacity it has.", CYAN),
            ("THE MIX", "Data selection: code, books, conversation, licensed data, synthetic examples, and filtering.", PURPLE),
            ("THE SEASONING", "Post-training shapes behavior, safety, tool use, reasoning style, and uncertainty.", AMBER),
        ],
        CYAN,
        "Models overlap in data and techniques, but no two labs use or publish exactly the same mixture.",
    )
    index += 1

    weights_artifact_slide(
        prs,
        index,
        "Connect directly to the previous slide: different recipes, data mixtures, and "
        "post-training choices produce different packaged models. This is one concrete model "
        "available locally rather than an abstract example.\n\n"
        "The left panel is a simplified snapshot of the local Ollama model listing. 27b "
        "indicates the approximate parameter count. q4_K_M identifies the quantization recipe "
        "used to store many weights with fewer bits. 17 GB is the resulting package size on disk.\n\n"
        "The package is not a folder containing 27 billion readable facts. It combines learned "
        "values and the information needed to run the model's architecture and tokenizer. "
        "Use q4_K_M as the handoff to the quantization slide that follows.",
    )
    index += 1

    compare_slide(
        prs,
        index,
        "Quantization",
        "Fewer digits, not fewer ideas",
        ("HIGHER PRECISION", "LARGER", ["Bigger model file", "More weight memory", "Reference-quality artifact"]),
        ("Q4-STYLE", "PRACTICAL", ["Much smaller weight storage", "Often faster on local hardware", "Some task-dependent quality loss"]),
        CYAN,
        "Quantization rounds many weights to fewer digits. Total RAM also includes context/KV cache and runtime buffers.",
    )
    index += 1

    token_slide(prs, index, "Tokens may be words, word pieces, punctuation, bytes, or special markers. The exact split depends on the tokenizer.")
    index += 1

    flow_slide(
        prs,
        index,
        "Inference",
        "What happens when you type a prompt",
        [
            ("TOKENS", "Text becomes token IDs.", PURPLE),
            ("RUN THE MATH", "The architecture applies the learned weights.", CYAN),
            ("SCORE", "Every possible next token receives a score.", AMBER),
            ("CHOOSE", "The system selects one token.", CORAL),
            ("LOOP", "Append it and repeat.", GREEN),
        ],
        CYAN,
        "This is inference: using fixed trained weights. The prompt and surrounding context "
        "are the changing inputs. Connect SCORE and CHOOSE back to the earlier vocabulary "
        "example, then use the next slide to summarize the repeated generation loop.",
        "One answer is built one token at a time.",
    )
    index += 1

    flow_slide(
        prs,
        index,
        "Generation",
        "How one answer gets built",
        [
            ("CONTEXT IN", "Take the conversation so far.", PURPLE),
            ("SCORES OUT", "Produce a score for every token.", CYAN),
            ("CHOOSE ONE", "Use decoding rules to select a token.", AMBER),
            ("FEED IT BACK", "Add the token to the context.", CORAL),
            ("STOP", "End when a stop condition is reached.", GREEN),
        ],
        CYAN,
        "Use the capital-of-France example. Paris wins, feeds back in, then punctuation and an end condition follow.",
        "Learned numbers + token vocabulary + repeated prediction.",
    )
    index += 1

    act_slide(
        prs,
        2,
        "Context and truth",
        "What the model can see, what it cannot know, and why fluent can still be wrong",
        PURPLE,
        "Transition from how tokens are generated to what information is available while generating them.",
    )
    index += 1

    context_window_slide(prs, index, "A base model call has no automatic memory of independent earlier calls. The product must supply relevant state again.")
    index += 1

    bob_slide(prs, index, "The Bob example makes compression concrete. Small personal details can disappear while project decisions survive.")
    index += 1

    three_cards_slide(
        prs,
        index,
        "Long context",
        "Bigger window does not mean smarter model",
        [
            (
                "CAPACITY",
                "A bigger notebook holds more pages.\n\nBut the prompt and answer must share the available space.",
                PURPLE,
            ),
            (
                "COST",
                "Every supplied page still has to be processed.\n\nMore context consumes memory and takes time to ingest.",
                AMBER,
            ),
            (
                "SIGNAL",
                "One critical sentence can hide on page 300.\n\nPresent does not always mean prominent.",
                BLUE,
            ),
        ],
        PURPLE,
        "Use a larger notebook as the analogy.\n\n"
        "CAPACITY: a bigger notebook can hold more pages. The context window is the maximum "
        "combined space for the supplied material and the answer.\n\n"
        "COST: someone still has to read those pages. More supplied context consumes memory "
        "and processing time even when much of it is not relevant to the current question.\n\n"
        "SIGNAL: the important sentence may still be present, but buried on page 300 among "
        "hundreds of pages of noise. The model may use prominent or nearby details more "
        "reliably than that buried detail.\n\n"
        "Contrast this with the previous Bob slide. After lossy compression, the detail may "
        "be completely gone. On this slide, the detail still fits but may not influence the "
        "answer strongly enough. The model window is an upper bound; the harness budget is "
        "the real session limit.",
    )
    index += 1

    current_information_slide(
        prs,
        index,
        "This is the bridge from context into the later harness and MCP sections.\n\n"
        "Start at the right: the model's trained weights are frozen for that model version. "
        "The model does not independently browse, call an API, or update those weights.\n\n"
        "The layout mirrors the headline: the stale model core is on the left and the fresh "
        "external source is at the far right.\n\n"
        "Trace the evidence path from right to left. A live source may contain current information. The harness "
        "retrieves it using a built-in tool or an external capability exposed through something "
        "such as an MCP server. The returned result is added to runtime context, and the same "
        "frozen model reads that evidence when producing its answer.\n\n"
        "A user pasting current information follows the same principle but skips the retrieval "
        "step. The information changes what the model can see now; it does not retrain the model.\n\n"
        "The answer is only as current and trustworthy as the supplied source. Act 3 explains "
        "the harness, tool-call loop, and MCP connection in detail.",
    )
    index += 1

    hallucination_slide(
        prs,
        index,
        "Use plain language: the model generates an answer, but it does not automatically perform "
        "an independent fact-check before showing that answer. Correct answers and hallucinations "
        "come out through the same token-by-token generation process.\n\n"
        "SAME PROCESS: the model does not switch into a visibly different 'hallucination mode.' "
        "The strongest learned continuation may be correct, incomplete, stale, or false.\n\n"
        "CHECKS VARY: the surrounding product or harness may add web search, citations, code "
        "execution, tests, critique, or another review step. Other products and requests simply "
        "show the first generated answer. Even when a check exists, its quality depends on the "
        "source and the check being performed.\n\n"
        "YOUR RULE: if no separate check occurred, treat the answer as unverified. For important "
        "claims, inspect a trusted source, run the code, test the result, or repeat the calculation. "
        "The next slide turns this rule into a short checklist.\n\n"
        "A model can be prompted to critique its own answer, but another generated opinion is not "
        "the same as independent evidence. Avoid that detour unless someone asks.\n\n"
        "Temperature can increase variety and sometimes increase error by giving weaker candidates "
        "more opportunity. It is not the root cause, and low temperature can still select a "
        "high-scoring answer that is confidently wrong.",
    )
    index += 1

    flow_slide(
        prs,
        index,
        "Verification",
        "Trust evidence, not presentation",
        [
            ("SOURCE", "Does the cited source exist?", CYAN),
            ("SUPPORT", "Does it actually support the claim?", PURPLE),
            ("EXECUTE", "Can the code or calculation run?", GREEN),
            ("TEST", "Does the result survive a real check?", AMBER),
            ("OWN IT", "A human remains accountable.", CORAL),
        ],
        PURPLE,
        "Generated reasoning and agreement between two models are clues, not proof.",
        "Verifiable artifacts beat confident assertions.",
    )
    index += 1

    act_slide(
        prs,
        3,
        "Giving the model hands",
        "Harnesses, tools, MCP, instructions, and reusable skills",
        CORAL,
        "Transition: the bare model predicts text. Agentic systems wrap that prediction loop in software that can act.",
    )
    index += 1

    statement_slide(
        prs,
        index,
        "The harness",
        "The model is the brain. The harness is the body.",
        "It assembles context, calls the model, enforces permissions, runs approved actions, and decides whether to continue.",
        CORAL,
        "Examples include Copilot CLI, Open WebUI, and other coding-agent applications.",
    )
    index += 1

    flow_slide(
        prs,
        index,
        "Agentic loop",
        "One common example: assemble -> call -> act -> repeat",
        [
            ("BUILD", "Assemble the information for this turn.", "rules + history + files", PURPLE),
            ("CALL", "Ask the model what should happen next.", '"run the tests"', CYAN),
            ("CHECK", "Apply permissions and guardrails.", "allow test; block deploy", AMBER),
            ("ACT", "Run an approved tool or operation.", "execute test command", CORAL),
            ("REPEAT", "Feed the result back and continue.", "test failure -> next turn", GREEN),
        ],
        CORAL,
        "This is a conceptual teaching pattern, not a required architecture or industry "
        "specification. Different harnesses are designed by different teams and may combine, "
        "split, reorder, parallelize, or omit these stages.\n\n"
        "Some systems make one model call and stop. Others use planners and executors, call "
        "multiple models, run tools in parallel, require human approval, or mix deterministic "
        "code with model decisions.\n\n"
        "The recurring idea is that software around the model assembles input, invokes the "
        "model, mediates allowed actions, returns results, and decides what happens next. "
        "An agent is the configured goal-seeking system; this repeated loop is one common way "
        "to implement it.\n\n"
        "The visible sub-examples depict one coding agent running tests. A research assistant, "
        "travel planner, or chatbot with web access would use different information and tools "
        "while still potentially following a similar high-level loop.",
        "Harnesses may combine, split, reorder, or skip these stages.",
        label="EXAMPLE SCENARIO: A CODING AGENT RUNNING TESTS",
    )
    index += 1

    tool_call_slide(prs, index, "Trace the request and response. The model does not execute the search itself.")
    index += 1

    bubbles_slide(
        prs,
        index,
        "Tool inventory",
        "A harness mounts capabilities",
        [
            ("FILES", "Search, view, edit, and patch", CYAN),
            ("SHELL", "Run builds, tests, and scripts", PURPLE),
            ("WEB", "Fetch current pages and APIs", AMBER),
            ("AGENTS", "Delegate specialized work", CORAL),
        ],
        CORAL,
        "This slide is the self-contained representative inventory. Exact tools vary by "
        "version, permissions, enabled features, and MCP servers, so do not present it as "
        "a permanent exhaustive list.",
        center="MODEL\nREQUESTS",
    )
    index += 1

    statement_slide(
        prs,
        index,
        "The file boundary",
        "The model sees only what the surrounding system supplies.",
        "Returned excerpts enter context. Unreturned file content stays outside. None of it becomes permanently trained into the weights.",
        CORAL,
        "This corrects the common idea that the model somehow has the local repository inside it.",
    )
    index += 1

    statement_slide(
        prs,
        index,
        "Why MCP?",
        "Stop rebuilding every connection pair-by-pair.",
        "MCP is a shared protocol for exposing capabilities to compatible AI hosts.",
        CORAL,
        "Use N x M custom integrations versus reusable hosts and servers as an idealized diagram, not a literal guarantee.",
    )
    index += 1

    mcp_slide(prs, index, "MCP separates the AI host/client connection from servers that expose tools, resources, and prompts.")
    index += 1

    stack_slide(
        prs,
        index,
        "Instructions",
        "The layers that steer one model",
        [
            ("SYSTEM", "Highest-level role, boundaries, and product rules.", CYAN),
            ("PROJECT", "Repository instructions such as AGENTS.md or CLAUDE.md.", PURPLE),
            ("HISTORY", "Prior conversation and returned tool results.", BLUE),
            ("USER", "The current request and constraints.", AMBER),
            ("UNTRUSTED DATA", "Web pages, files, and external text are data - not authority.", CORAL),
        ],
        CORAL,
        "Providers expose different role names and hierarchies. The stable idea is layered authority, not one flat competition for attention.",
    )
    index += 1

    three_cards_slide(
        prs,
        index,
        "System prompt",
        "Same weights. Different job.",
        [
            ("TUTOR", "Explain patiently. Ask guiding questions. Adapt to the learner.", CYAN),
            ("REVIEWER", "Find concrete defects. Demand evidence. Ignore style noise.", PURPLE),
            ("OPERATOR", "Use tools carefully. Respect permissions. Verify completion.", CORAL),
        ],
        CORAL,
        "A system prompt changes the frame without retraining the model. Some products expose it; others keep parts hidden.",
    )
    index += 1

    three_cards_slide(
        prs,
        index,
        "Skills",
        "Instructions, tools, and skills are different",
        [
            ("INSTRUCTION", "A rule loaded when its scope applies.", PURPLE),
            ("TOOL", "A capability the harness can execute.", CORAL),
            ("SKILL", "A reusable playbook for doing a specific job more consistently.", GREEN),
        ],
        CORAL,
        "A skill can package procedural guidance plus optional scripts and resources. It can be loaded automatically or explicitly.",
    )
    index += 1

    act_slide(
        prs,
        4,
        "Failure is part of the system",
        "Agentic tools need limits, boundaries, verification, and accountable humans",
        AMBER,
        "ACT 4 - TRUST + CONTROL\n"
        "This title is the workshop's editorial grouping, not an industry term. The act asks: "
        "now that the system can act, what makes it safe enough to trust with real work?\n"
        "The answer is not blind trust. It is control through permissions, stopping rules, "
        "validation, verification, and accountable humans.\n"
        "Failure modes are diagnostic patterns, not reasons to panic.",
    )
    index += 1

    failure_slide(prs, index, "Ask students which earlier concept explains each failure. This turns the entire workshop into a troubleshooting map.")
    index += 1

    compare_slide(
        prs,
        index,
        "Prompt injection",
        "Instructions and untrusted data are not the same",
        ("TRUSTED INSTRUCTIONS", "AUTHORITY", ["System and project rules", "Defined by the host", "Intended to control behavior"]),
        ("UNTRUSTED CONTENT", "DATA", ["Web pages and files", 'May say "ignore the rules"', "Must not gain authority"]),
        AMBER,
        "Models can still mistake data for an order. Defense requires permissions, isolation, validation, and approval - not wording alone.",
    )
    index += 1

    flow_slide(
        prs,
        index,
        "Responsible use",
        "Four habits worth keeping",
        [
            ("PROTECT", "Do not paste secrets, credentials, private records, or data you cannot share.", CYAN),
            ("VERIFY", "Check consequential facts, code, calculations, and citations.", PURPLE),
            ("LIMIT", "Give agents the minimum permissions and time they need.", CORAL),
            ("OWN IT", "A human remains accountable for the final decision and its consequences.", AMBER),
        ],
        AMBER,
        "These are durable habits for school, robotics, software work, and future AI products.",
        "Capability grows faster than trust. Keep the human accountable.",
    )
    index += 1

    recap_slide(prs, index, "Return to the model-versus-product thesis. The user experience is the complete system, not the weights alone.")
    index += 1

    closing_slide(
        prs,
        "Invite questions. Useful prompts: What surprised you? What would you trust today? What tool would you give an agent? What would you never let it do without approval?",
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output))
    except PermissionError:
        output = output.with_stem(f"{output.stem}_updated")
        prs.save(str(output))
    return output


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help=f"Output path (default: {DEFAULT_OUTPUT.name})",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()
    result = build_deck(args.output.resolve())
    print(f"Saved {result}")
